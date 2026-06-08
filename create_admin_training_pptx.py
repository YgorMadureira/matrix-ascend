"""
===============================================================================
MATRIX ASCEND — Guia Completo do Administrador (PowerPoint)
===============================================================================
Gera um PowerPoint profissional de treinamento para o perfil ADMIN.
Contém ~40 slides com screenshots anotados, setas, retângulos e instruções
passo a passo cobrindo todos os módulos da plataforma.

Uso: python create_admin_training_pptx.py
===============================================================================
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os, math

# ═══════════════════════════════════════════════════════════════
# PATHS
# ═══════════════════════════════════════════════════════════════
BRAIN = r"C:\Users\SPXBR38903\.gemini\antigravity\brain\da637c84-d029-4a63-a983-fe93d2312b58"
OUT   = r"C:\Users\SPXBR38903\Documents\Projetos AI\Matrix_Ascend_Admin_Training.pptx"

# ═══════════════════════════════════════════════════════════════
# SCREENSHOTS
# ═══════════════════════════════════════════════════════════════
SCREENSHOTS = {
    "dashboard":     os.path.join(BRAIN, "dashboard_page_1779130744731.png"),
    "collaborators": os.path.join(BRAIN, "collaborators_page_1779130754855.png"),
    "reports":       os.path.join(BRAIN, "reports_page_1779130766582.png"),
    "schedule":      os.path.join(BRAIN, "schedule_page_1779130777025.png"),
    "signatures":    os.path.join(BRAIN, "signatures_page_1779130786380.png"),
    "materials":     os.path.join(BRAIN, "materials_page_1779130797552.png"),
    "trainings":     os.path.join(BRAIN, "trainings_page_1779130806579.png"),
    "settings":      os.path.join(BRAIN, "settings_page_1779130815271.png"),
    "pts_requests":  os.path.join(BRAIN, "pts_requests_list_1779460485135.png"),
    "pts_form":      os.path.join(BRAIN, "pts_request_form_1779460598160.png"),
    "pts_approval":  os.path.join(BRAIN, "pts_approval_modal_1779460574438.png"),
}

# ═══════════════════════════════════════════════════════════════
# COLORS
# ═══════════════════════════════════════════════════════════════
DARK       = RGBColor(0x0D, 0x11, 0x17)
DARK2      = RGBColor(0x16, 0x1B, 0x22)
DARK3      = RGBColor(0x21, 0x26, 0x2D)
ACCENT     = RGBColor(0xEE, 0x4D, 0x2D)
ACCENT2    = RGBColor(0xFF, 0x6B, 0x4A)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY       = RGBColor(0x8B, 0x94, 0x9E)
LIGHT      = RGBColor(0xC9, 0xD1, 0xD9)
GREEN      = RGBColor(0x3F, 0xB9, 0x50)
BLUE       = RGBColor(0x58, 0xA6, 0xFF)
AMBER      = RGBColor(0xD2, 0x9C, 0x22)
RED        = RGBColor(0xE5, 0x53, 0x4B)
BORDER     = RGBColor(0x30, 0x36, 0x3D)
YELLOW_HL  = RGBColor(0xFF, 0xD7, 0x00)

# ═══════════════════════════════════════════════════════════════
# PRESENTATION SETUP
# ═══════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

# ═══════════════════════════════════════════════════════════════
# HELPER FUNCTIONS
# ═══════════════════════════════════════════════════════════════
def dark_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK

def add_accent_bar(slide, top=0, height=Inches(0.06)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, top, W, height)
    s.fill.solid(); s.fill.fore_color.rgb = ACCENT; s.line.fill.background()

def add_text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False,
             align=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.alignment = align
    r = p.runs[0]; r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold; r.font.name = font_name
    return txBox

def add_rich_text(slide, left, top, width, height, segments, size=14, align=PP_ALIGN.LEFT):
    """segments = [(text, color, bold), ...]"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    for text, color, bold in segments:
        r = p.add_run()
        r.text = text; r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold; r.font.name = "Segoe UI"
    return txBox

def add_multiline(slide, left, top, width, height, lines, size=13, color=LIGHT, line_spacing=1.5):
    """lines = list of strings"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.space_after = Pt(4)
        for r in p.runs:
            r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = "Segoe UI"

def add_bullet_list(slide, left, top, width, height, items, size=13, color=LIGHT, bullet="•"):
    """items = list of strings, each gets bullet prefix"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"  {bullet}  {item}"
        p.space_after = Pt(3)
        for r in p.runs:
            r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = "Segoe UI"
    return txBox

def add_numbered_list(slide, left, top, width, height, items, size=13, color=LIGHT):
    """items = list of strings, numbered automatically"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"  {i+1}.  {item}"
        p.space_after = Pt(4)
        for r in p.runs:
            r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = "Segoe UI"
    return txBox

def add_img(slide, path, left, top, width, height=None):
    if os.path.exists(path):
        if height:
            return slide.shapes.add_picture(path, left, top, width, height)
        else:
            return slide.shapes.add_picture(path, left, top, width=width)
    else:
        print(f"  ⚠ Screenshot não encontrada: {path}")
        return None

def add_screenshot_with_border(slide, key, left, top, width, height, title=None):
    """Adds a screenshot with a dark rounded border"""
    # Border background
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left - Inches(0.1), top - Inches(0.1),
                                width + Inches(0.2), height + Inches(0.2))
    s.fill.solid(); s.fill.fore_color.rgb = DARK3
    s.line.color.rgb = BORDER
    # Image
    img_path = SCREENSHOTS.get(key, "")
    add_img(slide, img_path, left, top, width, height)

def add_callout_rect(slide, left, top, width, height, color=ACCENT, line_width=Pt(3)):
    """Draws a highlighted rectangle border (annotation)"""
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.background()
    s.line.color.rgb = color
    s.line.width = line_width
    return s

def add_callout_arrow(slide, text, left, top, width, height, direction="right", color=ACCENT):
    """Adds a labeled arrow callout"""
    # Arrow shape
    if direction == "right":
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, Inches(0.8), Inches(0.35))
    elif direction == "left":
        arrow = slide.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, left, top, Inches(0.8), Inches(0.35))
    elif direction == "down":
        arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, Inches(0.35), Inches(0.6))
    elif direction == "up":
        arrow = slide.shapes.add_shape(MSO_SHAPE.UP_ARROW, left, top, Inches(0.35), Inches(0.6))
    else:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, Inches(0.8), Inches(0.35))
    
    arrow.fill.solid(); arrow.fill.fore_color.rgb = color; arrow.line.fill.background()
    # Label
    add_text(slide, left + Inches(0.9) if direction in ("right",) else left - Inches(3),
             top - Inches(0.05), width, height, text, 11, color, True)

def add_step_badge(slide, left, top, number, size=Inches(0.5)):
    """Adds a numbered circular badge"""
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    s.fill.solid(); s.fill.fore_color.rgb = ACCENT; s.line.fill.background()
    tf = s.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.text = str(number)
    r = p.runs[0]; r.font.size = Pt(18); r.font.color.rgb = WHITE; r.font.bold = True; r.font.name = "Segoe UI"
    tf.paragraphs[0].space_before = Pt(0)
    try:
        s.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        bodyPr = s.text_frame._txBody.find(qn('a:bodyPr'))
        bodyPr.set('anchor', 'ctr')
    except: pass

def add_section_header(slide, icon, title, subtitle):
    """Standard section header"""
    add_text(slide, Inches(0.8), Inches(0.3), Inches(0.6), Inches(0.5), icon, 28, ACCENT)
    add_text(slide, Inches(1.5), Inches(0.3), Inches(10), Inches(0.7), title, 30, ACCENT, True)
    add_text(slide, Inches(1.5), Inches(0.9), Inches(10), Inches(0.4), subtitle, 13, GRAY)

def add_info_card(slide, left, top, width, height, title, lines, icon="ℹ️"):
    """Adds an info card with title and bullet lines"""
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = DARK2; s.line.color.rgb = BORDER
    add_text(slide, left + Inches(0.3), top + Inches(0.15), width - Inches(0.4), Inches(0.3),
             f"{icon}  {title}", 12, ACCENT, True)
    for i, line in enumerate(lines):
        add_text(slide, left + Inches(0.3), top + Inches(0.5 + i * 0.3), width - Inches(0.6), Inches(0.25),
                 f"• {line}", 11, LIGHT)

def new_slide():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(sl); add_accent_bar(sl)
    return sl

def add_slide_number(slide, num, total):
    add_text(slide, Inches(12), Inches(7.1), Inches(1), Inches(0.3),
             f"{num}/{total}", 9, GRAY, False, PP_ALIGN.RIGHT)

TOTAL_SLIDES = 42

print("🎯 Gerando PowerPoint de Treinamento ADMIN — Matrix Ascend...")
print(f"   Destino: {OUT}")
print(f"   Total de slides: {TOTAL_SLIDES}")
print()

slide_num = 0

# ╔═══════════════════════════════════════════════════════════════╗
# ║  SLIDE 1 — CAPA                                              ║
# ╚═══════════════════════════════════════════════════════════════╝
slide_num += 1
print(f"  [{slide_num:02d}] Capa")
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl); add_accent_bar(sl)

s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1), W, Inches(5.5))
s.fill.solid(); s.fill.fore_color.rgb = DARK2; s.line.fill.background()

add_text(sl, Inches(1), Inches(1.5), Inches(11), Inches(1), "MATRIX ASCEND", 60, ACCENT, True)
add_text(sl, Inches(1), Inches(2.5), Inches(10), Inches(0.8), "Guia Completo do Administrador", 32, WHITE, True)
add_text(sl, Inches(1), Inches(3.3), Inches(10), Inches(0.6), "Manual de Treinamento — Perfil ADMIN", 20, LIGHT)
add_text(sl, Inches(1), Inches(4.0), Inches(10), Inches(0.5), "Shopee Logistics  •  SPX BR  •  Portal de Treinamentos", 14, GRAY)

s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(4.8), Inches(2.5), Inches(0.04))
s.fill.solid(); s.fill.fore_color.rgb = ACCENT; s.line.fill.background()

add_text(sl, Inches(1), Inches(5.2), Inches(8), Inches(0.5),
         "Este documento contém instruções completas e passo a passo para todas as", 12, GRAY)
add_text(sl, Inches(1), Inches(5.5), Inches(8), Inches(0.5),
         "funcionalidades da plataforma no perfil Administrador.", 12, GRAY)

add_text(sl, Inches(1), Inches(6.5), Inches(4), Inches(0.4), "v2.4  •  Sistema Online e Operacional", 12, GREEN, True)
add_text(sl, Inches(8), Inches(6.5), Inches(4), Inches(0.4), "© 2026 Matrix Ascend — SPX BR", 12, GRAY, False, PP_ALIGN.RIGHT)
add_slide_number(sl, slide_num, TOTAL_SLIDES)

# ╔═══════════════════════════════════════════════════════════════╗
# ║  SLIDE 2 — SUMÁRIO                                           ║
# ╚═══════════════════════════════════════════════════════════════╝
slide_num += 1
print(f"  [{slide_num:02d}] Sumário")
sl = new_slide()

add_text(sl, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7), "SUMÁRIO DO TREINAMENTO", 30, ACCENT, True)
add_text(sl, Inches(0.8), Inches(0.9), Inches(10), Inches(0.4), "Todos os tópicos cobertos neste material", 13, GRAY)

chapters = [
    ("01", "Visão Geral da Plataforma", "Navegação, perfis de acesso e primeiro login"),
    ("02", "Dashboard", "Painel de KPIs e desempenho por macro-setor"),
    ("03", "Configurações — Usuários", "Como criar, editar e remover usuários da plataforma"),
    ("04", "Configurações — Instrutores", "Como cadastrar e gerenciar instrutores"),
    ("05", "Configurações — Avaliações", "Banco de questões para provas de certificação"),
    ("06", "Configurações — Sincronização", "Sobre a sincronização da base ABS"),
    ("07", "Materiais", "Gestão de pastas e links do Google Drive + QR Code"),
    ("08", "Colaboradores", "Base ativa, onboarding, importação CSV e badges"),
    ("09", "Treinamentos EAD", "Capacitação: vídeo, prova e assinatura digital"),
    ("10", "Agenda", "Calendário, slots, solicitações e aprovações"),
    ("11", "SOCs — Unidades", "Cadastro e gestão de unidades operacionais"),
    ("12", "Assinaturas e Relatórios", "Auditoria digital e relatórios de certificação"),
]

for i, (num, title, desc) in enumerate(chapters):
    col = i // 6
    row = i % 6
    x = Inches(0.8 + col * 6)
    y = Inches(1.5 + row * 0.92)
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.5), Inches(0.8))
    s.fill.solid(); s.fill.fore_color.rgb = DARK2; s.line.fill.background()
    add_text(sl, x + Inches(0.2), y + Inches(0.1), Inches(0.5), Inches(0.5), num, 16, ACCENT, True)
    add_text(sl, x + Inches(0.8), y + Inches(0.08), Inches(4.5), Inches(0.35), title, 13, WHITE, True)
    add_text(sl, x + Inches(0.8), y + Inches(0.42), Inches(4.5), Inches(0.3), desc, 10, GRAY)

add_slide_number(sl, slide_num, TOTAL_SLIDES)

# ╔═══════════════════════════════════════════════════════════════╗
# ║  SLIDE 3 — VISÃO GERAL DA PLATAFORMA                        ║
# ╚═══════════════════════════════════════════════════════════════╝
slide_num += 1
print(f"  [{slide_num:02d}] Visão Geral da Plataforma")
sl = new_slide()
add_section_header(sl, "🖥️", "VISÃO GERAL DA PLATAFORMA", "Conheça a estrutura de navegação e os módulos disponíveis")

# Screenshot do dashboard com anotações
add_screenshot_with_border(sl, "dashboard", Inches(0.5), Inches(1.5), Inches(7.5), Inches(5.2))

# Anotações sobre a sidebar
add_callout_rect(sl, Inches(0.5), Inches(1.5), Inches(1.8), Inches(5.2), YELLOW_HL, Pt(3))
add_text(sl, Inches(8.5), Inches(1.5), Inches(4.5), Inches(0.4), "BARRA LATERAL (SIDEBAR)", 14, ACCENT, True)

sidebar_items = [
    "📊  Dashboard — Painel de KPIs",
    "📁  Materiais — Links Google Drive",
    "👥  Colaboradores — Base de funcionários",
    "📈  Relatórios — Gráficos e matrizes",
    "📅  Agenda — Calendário de treinamentos",
    "✍️  Assinaturas — Registros digitais",
    "🏢  SOCs — Unidades operacionais",
    "🎓  Treinamentos — Capacitação EAD",
    "⚙️  Configurações — EXCLUSIVO ADMIN",
]
for i, item in enumerate(sidebar_items):
    color = ACCENT if "ADMIN" in item else LIGHT
    bold = "ADMIN" in item
    add_text(sl, Inches(8.5), Inches(2.0 + i * 0.42), Inches(4.5), Inches(0.35), item, 12, color, bold)

add_info_card(sl, Inches(8.5), Inches(6.0), Inches(4.3), Inches(0.8),
              "DICA IMPORTANTE", ["O menu 'Configurações' só aparece para ADMIN"])
add_slide_number(sl, slide_num, TOTAL_SLIDES)

# ╔═══════════════════════════════════════════════════════════════╗
# ║  SLIDE 4 — BARRA SUPERIOR E HEADER                          ║
# ╚═══════════════════════════════════════════════════════════════╝
slide_num += 1
print(f"  [{slide_num:02d}] Header e Navegação")
sl = new_slide()
add_section_header(sl, "🔝", "HEADER E NAVEGAÇÃO", "Entendendo o cabeçalho e as informações de perfil")

# Screenshot focada no topo
add_screenshot_with_border(sl, "dashboard", Inches(0.5), Inches(1.5), Inches(12), Inches(1.5))

# Anotações do header
add_callout_rect(sl, Inches(0.5), Inches(1.5), Inches(2.8), Inches(1.0), YELLOW_HL, Pt(3))
add_text(sl, Inches(0.5), Inches(3.2), Inches(3), Inches(0.3), "1️⃣  Logo e nome do portal", 12, LIGHT)

add_callout_rect(sl, Inches(9.5), Inches(1.5), Inches(3), Inches(1.0), GREEN, Pt(3))
add_text(sl, Inches(9.5), Inches(3.2), Inches(3), Inches(0.3), "2️⃣  Nome do usuário + perfil ADMIN", 12, LIGHT)

add_callout_rect(sl, Inches(12.0), Inches(1.5), Inches(0.8), Inches(1.0), BLUE, Pt(3))
add_text(sl, Inches(10.5), Inches(3.6), Inches(3), Inches(0.3), "3️⃣  Botão de logout (sair)", 12, LIGHT)

# Explicação dos perfis
add_text(sl, Inches(0.8), Inches(4.3), Inches(10), Inches(0.5), "PERFIS DE ACESSO DISPONÍVEIS:", 16, ACCENT, True)

profiles = [
    ("ADMIN", "Acesso total — todos os módulos + Configurações", "🔴"),
    ("USUÁRIO", "Dashboard, Materiais, Colaboradores, Relatórios, SOCs, Treinamentos, Assinaturas", "🔵"),
    ("LÍDER", "Dashboard, Colaboradores (só seu time), Relatórios, Treinamentos, Agenda", "🟡"),
    ("BPO", "Colaboradores (aba Onboarding) — NÃO deleta registros", "🟢"),
    ("PCP", "Agenda — gerencia inscrições de treinamentos", "🟣"),
]

for i, (name, desc, icon) in enumerate(profiles):
    y = Inches(4.9 + i * 0.5)
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(11.5), Inches(0.42))
    s.fill.solid(); s.fill.fore_color.rgb = DARK2; s.line.fill.background()
    add_text(sl, Inches(1.0), y + Inches(0.05), Inches(0.3), Inches(0.3), icon, 14, WHITE)
    add_text(sl, Inches(1.5), y + Inches(0.05), Inches(1.5), Inches(0.3), name, 12, WHITE, True)
    add_text(sl, Inches(3.2), y + Inches(0.05), Inches(9), Inches(0.3), desc, 11, GRAY)

add_slide_number(sl, slide_num, TOTAL_SLIDES)

# ╔═══════════════════════════════════════════════════════════════╗
# ║  SLIDE 5 — PRIMEIRO ACESSO / LOGIN                          ║
# ╚═══════════════════════════════════════════════════════════════╝
slide_num += 1
print(f"  [{slide_num:02d}] Login e Primeiro Acesso")
sl = new_slide()
add_section_header(sl, "🔐", "LOGIN E PRIMEIRO ACESSO", "Como acessar a plataforma e o fluxo de troca de senha obrigatória")

# Coluna esquerda - Fluxo de login
add_text(sl, Inches(0.8), Inches(1.5), Inches(5), Inches(0.5), "FLUXO DE PRIMEIRO ACESSO:", 16, WHITE, True)

steps = [
    "Acesse a URL da plataforma no navegador",
    "Informe o e-mail corporativo cadastrado",
    "Insira a senha provisória fornecida pelo ADMIN",
    "No primeiro acesso, será exibida a tela de TROCA DE SENHA OBRIGATÓRIA",
    "Defina uma nova senha (mínimo 6 caracteres)",
    "Confirme a nova senha e clique em 'Redefinir Senha'",
    "Pronto! Você será redirecionado ao Dashboard",
]

for i, step in enumerate(steps):
    y = Inches(2.1 + i * 0.65)
    add_step_badge(sl, Inches(0.8), y, i + 1, Inches(0.45))
    add_text(sl, Inches(1.5), y + Inches(0.05), Inches(5), Inches(0.45), step, 12, LIGHT)

# Coluna direita - Informações importantes
add_info_card(sl, Inches(7), Inches(1.5), Inches(5.5), Inches(1.8),
              "SOBRE A SENHA PROVISÓRIA",
              [
                  "O ADMIN define a senha ao criar o usuário",
                  "O sistema força a troca no primeiro login",
                  "A flag 'must_change_password' é automática",
                  "Após trocar, o acesso é imediato",
              ])

add_info_card(sl, Inches(7), Inches(3.6), Inches(5.5), Inches(1.5),
              "IMPORTANTE PARA O ADMIN",
              [
                  "Ao criar um usuário, informe a senha provisória",
                  "Oriente o novo PTS a trocar no primeiro acesso",
                  "A senha deve ter no mínimo 6 caracteres",
              ])

add_info_card(sl, Inches(7), Inches(5.4), Inches(5.5), Inches(1.2),
              "URL DE ACESSO",
              [
                  "Produção: matrix-ascend.vercel.app",
                  "Funciona em desktop, tablet e smartphone",
              ])

add_slide_number(sl, slide_num, TOTAL_SLIDES)

# ╔═══════════════════════════════════════════════════════════════╗
# ║  SLIDE 6 — DASHBOARD                                         ║
# ╚═══════════════════════════════════════════════════════════════╝
slide_num += 1
print(f"  [{slide_num:02d}] Dashboard")
sl = new_slide()
add_section_header(sl, "📊", "DASHBOARD", "Visão geral dos indicadores — a primeira tela ao entrar no sistema")

add_screenshot_with_border(sl, "dashboard", Inches(0.5), Inches(1.4), Inches(7.8), Inches(5.5))

# Anotações nos cards de KPI
add_callout_rect(sl, Inches(2.5), Inches(3.2), Inches(5.5), Inches(1.2), YELLOW_HL, Pt(2))
add_text(sl, Inches(8.5), Inches(1.5), Inches(4.5), Inches(0.4), "CARDS DE KPI:", 14, ACCENT, True)
kpis = [
    ("👥 MEU TIME", "Total de colaboradores da SOC"),
    ("📊 % TREINADOS", "Percentual de certificação global"),
    ("📁 MATERIAIS", "Total de materiais cadastrados"),
    ("✅ TREINADOS", "Número absoluto de treinados"),
]
for i, (name, desc) in enumerate(kpis):
    y = Inches(2.0 + i * 0.55)
    add_text(sl, Inches(8.5), y, Inches(2), Inches(0.3), name, 11, WHITE, True)
    add_text(sl, Inches(10.5), y, Inches(3), Inches(0.3), desc, 10, GRAY)

# Anotação macro-setor
add_callout_rect(sl, Inches(2.5), Inches(4.8), Inches(5.5), Inches(2.0), GREEN, Pt(2))
add_text(sl, Inches(8.5), Inches(4.5), Inches(4.5), Inches(0.4), "DESEMPENHO POR MACRO-SETOR:", 14, ACCENT, True)
sectors = ["Recebimento", "Processamento", "Expedição", "Tratativas", "HSE", "People"]
for i, s in enumerate(sectors):
    add_text(sl, Inches(8.5), Inches(5.0 + i * 0.35), Inches(4), Inches(0.3),
             f"  •  {s} — % de certificação e total", 10, LIGHT)

add_slide_number(sl, slide_num, TOTAL_SLIDES)

# ╔═══════════════════════════════════════════════════════════════╗
# ║  SLIDE 7 — CONFIGURAÇÕES: VISÃO GERAL                       ║
# ╚═══════════════════════════════════════════════════════════════╝
slide_num += 1
print(f"  [{slide_num:02d}] Configurações - Visão Geral")
sl = new_slide()
add_section_header(sl, "⚙️", "CONFIGURAÇÕES — VISÃO GERAL", "Central de controle exclusiva do perfil ADMIN")

add_screenshot_with_border(sl, "settings", Inches(0.5), Inches(1.4), Inches(7.8), Inches(5.5))

# Anotações
add_callout_rect(sl, Inches(2.5), Inches(3.3), Inches(2.6), Inches(3.4), YELLOW_HL, Pt(3))
add_text(sl, Inches(5.3), Inches(3.3), Inches(0.8), Inches(0.35), "→", 20, YELLOW_HL, True)
add_text(sl, Inches(6.0), Inches(3.3), Inches(2), Inches(0.35), "Painel de Usuários", 12, YELLOW_HL, True)

add_callout_rect(sl, Inches(5.5), Inches(3.3), Inches(2.6), Inches(3.4), GREEN, Pt(3))
add_text(sl, Inches(6.0), Inches(3.7), Inches(2), Inches(0.35), "Painel de Instrutores", 12, GREEN, True)

add_text(sl, Inches(8.8), Inches(1.5), Inches(4.2), Inches(0.4), "SEÇÕES DISPONÍVEIS:", 14, ACCENT, True)

sections = [
    ("👤 USUÁRIOS", "Crie, edite e remova acessos"),
    ("🎓 INSTRUTORES", "Cadastre instrutores por SOC"),
    ("📝 AVALIAÇÕES", "Configure provas de certificação"),
    ("🔄 SINCRONIZAÇÃO", "Horário de atualização da base"),
]

for i, (name, desc) in enumerate(sections):
    y = Inches(2.1 + i * 0.7)
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), y, Inches(4.2), Inches(0.55))
    s.fill.solid(); s.fill.fore_color.rgb = DARK2; s.line.fill.background()
    add_text(sl, Inches(9.0), y + Inches(0.05), Inches(3.8), Inches(0.25), name, 12, WHITE, True)
    add_text(sl, Inches(9.0), y + Inches(0.28), Inches(3.8), Inches(0.25), desc, 10, GRAY)

add_info_card(sl, Inches(8.8), Inches(5.2), Inches(4.2), Inches(1.2),
              "ACESSO RESTRITO", [
                  "Apenas perfis ADMIN veem este menu",
                  "Se não aparecer, verifique o perfil do usuário",
                  "O link fica na seção 'ADMINISTRAÇÃO' da sidebar",
              ])
add_slide_number(sl, slide_num, TOTAL_SLIDES)

# ╔═══════════════════════════════════════════════════════════════╗
# ║  SLIDE 8 — CRIAR NOVO USUÁRIO (passo a passo)               ║
# ╚═══════════════════════════════════════════════════════════════╝
slide_num += 1
print(f"  [{slide_num:02d}] Criar Novo Usuário - Passo a Passo")
sl = new_slide()
add_section_header(sl, "👤", "CRIAR NOVO USUÁRIO", "Passo a passo para cadastrar um novo PTS ou acesso na plataforma")

# Steps
add_text(sl, Inches(0.8), Inches(1.5), Inches(12), Inches(0.5), "PASSO A PASSO:", 16, WHITE, True)

steps_data = [
    ("1", "Acesse 'Configurações' no menu lateral", "Clique em 'Configurações' na seção ADMINISTRAÇÃO da sidebar"),
    ("2", "Clique no ícone de '+ Usuário'", "No painel USUÁRIOS, clique no ícone de pessoa com '+' no canto superior direito"),
    ("3", "Preencha o formulário 'Novo Usuário'", "Um modal será aberto com todos os campos necessários (veja slide seguinte)"),
    ("4", "Selecione o Nível de Acesso", "Escolha entre: Usuário Comum, Líder Operacional, BPO, PCP ou Administrador"),
    ("5", "Selecione a Unidade Base (SOC)", "Escolha a SOC correspondente ao PTS (Ex: SP6, SP1, etc.)"),
    ("6", "Defina uma senha provisória", "Mínimo 6 caracteres — o usuário será obrigado a trocar no primeiro acesso"),
    ("7", "Confirme a senha", "Repita a mesma senha — indicadores visuais: ✅ verde = coincidem, ❌ vermelho = não coincidem"),
    ("8", "Clique em 'CADASTRAR E SINCRONIZAR'", "O usuário será criado imediatamente e poderá fazer login"),
]

for i, (num, title, desc) in enumerate(steps_data):
    y = Inches(2.0 + i * 0.65)
    add_step_badge(sl, Inches(0.8), y, int(num), Inches(0.4))
    add_text(sl, Inches(1.4), y + Inches(0.02), Inches(5), Inches(0.3), title, 13, WHITE, True)
    add_text(sl, Inches(6.5), y + Inches(0.02), Inches(6), Inches(0.3), desc, 11, GRAY)

add_slide_number(sl, slide_num, TOTAL_SLIDES)

# ╔═══════════════════════════════════════════════════════════════╗
# ║  SLIDE 9 — FORMULÁRIO DE NOVO USUÁRIO (campos)              ║
# ╚═══════════════════════════════════════════════════════════════╝
slide_num += 1
print(f"  [{slide_num:02d}] Formulário de Novo Usuário")
sl = new_slide()
add_section_header(sl, "📝", "FORMULÁRIO — NOVO USUÁRIO", "Detalhes de cada campo do modal de criação de usuário")

# Mock form visual
form_fields = [
    ("NOME COMPLETO", "Nome completo do PTS ou colaborador", "Ex: João Silva Santos", "Texto obrigatório"),
    ("E-MAIL CORPORATIVO", "E-mail que será usado para login", "Ex: joao.silva@shopee.com", "Obrigatório, único"),
    ("SENHA PROVISÓRIA", "Senha inicial para o primeiro acesso", "Mínimo 6 caracteres", "Será trocada no 1º login"),
    ("CONFIRMAR SENHA", "Repetição da senha provisória", "Deve ser igual à senha acima", "Indicadores ✅/❌"),
    ("NÍVEL DE ACESSO", "Perfil de permissões do usuário", "Dropdown com 5 opções", "Define o que o usuário vê"),
    ("UNIDADE BASE (SOC)", "SOC onde o usuário atua", "Dropdown com SOCs cadastradas", "Filtra dados por unidade"),
]

for i, (label, desc, example, obs) in enumerate(form_fields):
    y = Inches(1.5 + i * 0.85)
    # Campo visual
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(5.5), Inches(0.72))
    s.fill.solid(); s.fill.fore_color.rgb = DARK2; s.line.color.rgb = BORDER
    add_text(sl, Inches(1.0), y + Inches(0.05), Inches(2.5), Inches(0.25), label, 10, ACCENT, True)
    add_text(sl, Inches(1.0), y + Inches(0.3), Inches(5), Inches(0.25), example, 11, GRAY)
    # Descrição
    add_text(sl, Inches(6.8), y + Inches(0.05), Inches(3), Inches(0.25), desc, 11, LIGHT)
    add_text(sl, Inches(10.3), y + Inches(0.05), Inches(3), Inches(0.25), obs, 10, AMBER, True)

# Opções do dropdown Nível de Acesso
add_text(sl, Inches(0.8), Inches(6.8), Inches(12), Inches(0.3),
         "💡 OPÇÕES DO NÍVEL DE ACESSO:", 12, ACCENT, True)
options = ["Usuário Comum (user) — Consulta Limitada",
           "Líder Operacional (lider) — Gestão do Time",
           "BPO Onboarding (bpo) — Aba Onboarding",
           "Administrador Geral (admin) — Acesso Total",
           "PCP Agenda (pcp) — Gerencia Agenda"]
add_text(sl, Inches(0.8), Inches(7.1), Inches(12), Inches(0.3),
         "   •   ".join(options), 9, LIGHT)

add_slide_number(sl, slide_num, TOTAL_SLIDES)

# ╔═══════════════════════════════════════════════════════════════╗
# ║  SLIDE 10 — EDITAR USUÁRIO                                   ║
# ╚═══════════════════════════════════════════════════════════════╝
slide_num += 1
print(f"  [{slide_num:02d}] Editar Usuário")
sl = new_slide()
add_section_header(sl, "✏️", "EDITAR PERFIL DE USUÁRIO", "Como alterar nome, perfil de acesso e unidade de um usuário existente")

# Steps
steps = [
    "Na tela de Configurações, localize o usuário no painel USUÁRIOS",
    "Passe o mouse sobre o nome do usuário — aparecem os ícones de EDITAR (✏️) e EXCLUIR (🗑️)",
    "Clique no ícone de EDITAR (lápis) para abrir o modal 'Editar Perfil'",
    "Altere os campos desejados: Nome, Regra de Acesso, Unidade Base (SOC)",
    "Para Líderes: preencha a 'Chave de Identificação' (nome EXATO como aparece na planilha)",
    "Clique em 'SALVAR' para confirmar as alterações",
]

for i, step in enumerate(steps):
    y = Inches(1.5 + i * 0.7)
    add_step_badge(sl, Inches(0.8), y, i + 1, Inches(0.4))
    add_text(sl, Inches(1.5), y + Inches(0.03), Inches(6), Inches(0.5), step, 13, LIGHT)

# Info cards
add_info_card(sl, Inches(8), Inches(1.5), Inches(4.8), Inches(2.0),
              "CAMPOS EDITÁVEIS", [
                  "Nome — Nome completo do usuário",
                  "Regra de Acesso — admin/user/lider/bpo/pcp",
                  "Unidade Base (SOC) — SOC vinculada",
                  "Chave do Líder — Só para perfil 'lider'",
              ])

add_info_card(sl, Inches(8), Inches(3.8), Inches(4.8), Inches(1.8),
              "CHAVE DE IDENTIFICAÇÃO (LÍDER)", [
                  "Aparece APENAS quando perfil = Líder",
                  "Deve ser o nome EXATO da coluna 'Líder'",
                  "Permite que o líder veja apenas seu time",
                  "Ex: RICARDO MARTINS (conforme planilha)",
              ])

add_info_card(sl, Inches(8), Inches(5.9), Inches(4.8), Inches(1.2),
              "⚠️ EXCLUIR USUÁRIO", [
                  "Clique no ícone de lixeira (🗑️)",
                  "Confirme na caixa de diálogo",
                  "A EXCLUSÃO é permanente!",
              ])

add_slide_number(sl, slide_num, TOTAL_SLIDES)

# ╔═══════════════════════════════════════════════════════════════╗
# ║  SLIDE 11 — INSTRUTORES                                      ║
# ╚═══════════════════════════════════════════════════════════════╝
slide_num += 1
print(f"  [{slide_num:02d}] Cadastrar Instrutores")
sl = new_slide()
add_section_header(sl, "🎓", "CADASTRAR INSTRUTORES", "Como adicionar e gerenciar os instrutores de treinamento")

add_screenshot_with_border(sl, "settings", Inches(0.5), Inches(1.4), Inches(6.5), Inches(4.5))

# Anotação no painel de instrutores
add_callout_rect(sl, Inches(5.0), Inches(3.0), Inches(2.0), Inches(2.8), GREEN, Pt(3))

add_text(sl, Inches(7.5), Inches(1.5), Inches(5.3), Inches(0.5), "PASSO A PASSO:", 14, ACCENT, True)

steps = [
    "Na tela de Configurações, localize o painel 'INSTRUTORES'",
    "Clique no botão '+' no canto superior direito do painel",
    "No modal 'Novo Instrutor', preencha:",
    "   → Nome Completo (ex: Rodrigo Souza)",
    "   → Unidade Base (SOC) — selecione no dropdown",
    "Clique em 'CONFIRMAR INSTRUCTOR'",
    "O instrutor aparecerá na lista com o ícone 🎓 e a base",
]
for i, step in enumerate(steps):
    add_text(sl, Inches(7.5), Inches(2.0 + i * 0.45), Inches(5.3), Inches(0.35),
             f"  {i+1}.  {step}" if not step.startswith("   ") else step, 11, LIGHT)

add_info_card(sl, Inches(7.5), Inches(5.4), Inches(5.3), Inches(1.5),
              "GERENCIAR INSTRUTORES", [
                  "Para EXCLUIR: passe o mouse e clique 🗑️",
                  "Instrutores aparecem na Agenda ao criar slots",
                  "Cada instrutor é vinculado a uma SOC",
                  "Não há limite de instrutores por SOC",
              ])

add_slide_number(sl, slide_num, TOTAL_SLIDES)

# ╔═══════════════════════════════════════════════════════════════╗
# ║  SLIDE 12 — AVALIAÇÕES E PROVAS                              ║
# ╚═══════════════════════════════════════════════════════════════╝
slide_num += 1
print(f"  [{slide_num:02d}] Avaliações e Provas")
sl = new_slide()
add_section_header(sl, "📝", "AVALIAÇÕES E PROVAS", "Configure o banco de questões para certificação dos colaboradores")

add_text(sl, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
         "COMO FUNCIONA O SISTEMA DE PROVAS:", 16, WHITE, True)

# Visual do fluxo
flow_items = [
    ("1️⃣", "ADMIN cadastra as questões", "Na aba Avaliações, selecione o treinamento e adicione questões"),
    ("2️⃣", "Colaborador assiste ao vídeo", "Precisa assistir 60 segundos mínimos do conteúdo EAD"),
    ("3️⃣", "Colaborador faz a prova", "Múltipla escolha (A, B, C, D) — precisa de 90% para aprovar"),
    ("4️⃣", "Até 5 tentativas", "Após 5 reprovações, deve reassistir o vídeo completo"),
    ("5️⃣", "Assinatura digital", "Aprovado → assina digitalmente no canvas → certificação registrada"),
]

for i, (icon, title, desc) in enumerate(flow_items):
    y = Inches(2.1 + i * 0.7)
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(11.5), Inches(0.58))
    s.fill.solid(); s.fill.fore_color.rgb = DARK2; s.line.fill.background()
    add_text(sl, Inches(1.0), y + Inches(0.08), Inches(0.5), Inches(0.4), icon, 16, WHITE)
    add_text(sl, Inches(1.7), y + Inches(0.05), Inches(3.5), Inches(0.4), title, 13, WHITE, True)
    add_text(sl, Inches(5.3), y + Inches(0.05), Inches(7), Inches(0.4), desc, 11, GRAY)

# Instruções para criar questão
add_text(sl, Inches(0.8), Inches(5.8), Inches(11), Inches(0.4), "COMO ADICIONAR QUESTÕES:", 14, ACCENT, True)
add_numbered_list(sl, Inches(0.8), Inches(6.2), Inches(11), Inches(1.5), [
    "Em Configurações → seção 'Avaliações e Provas', clique no card do treinamento desejado",
    "Clique em 'ADICIONAR QUESTÃO' — preencha: Enunciado, Opções A/B/C/D, Resposta Correta, Ordem",
    "Clique em 'FINALIZAR E SALVAR QUESTÃO' — repita para cada questão da prova",
], 11, LIGHT)

add_slide_number(sl, slide_num, TOTAL_SLIDES)

# ╔═══════════════════════════════════════════════════════════════╗
# ║  SLIDE 13 — DETALHES DA QUESTÃO                              ║
# ╚═══════════════════════════════════════════════════════════════╝
slide_num += 1
print(f"  [{slide_num:02d}] Detalhes da Questão")
sl = new_slide()
add_section_header(sl, "❓", "FORMULÁRIO DE QUESTÃO", "Detalhes de cada campo ao criar ou editar uma questão")

fields = [
    ("ORDEM", "Número sequencial da questão na prova", "Ex: 1, 2, 3...", "Numeral"),
    ("ENUNCIADO", "Texto da pergunta que o colaborador verá", "Ex: Qual é o EPI obrigatório no setor...?", "Texto"),
    ("OPÇÃO A", "Primeira alternativa de resposta", "Ex: Capacete", "Texto"),
    ("OPÇÃO B", "Segunda alternativa de resposta", "Ex: Luvas", "Texto"),
    ("OPÇÃO C", "Terceira alternativa de resposta", "Ex: Óculos", "Texto"),
    ("OPÇÃO D", "Quarta alternativa de resposta", "Ex: Protetor auricular", "Texto"),
    ("RESPOSTA CORRETA", "Qual das 4 opções é a certa", "Selecione A, B, C ou D", "Botão"),
]

for i, (label, desc, example, tipo) in enumerate(fields):
    y = Inches(1.5 + i * 0.72)
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(11.5), Inches(0.62))
    s.fill.solid(); s.fill.fore_color.rgb = DARK2; s.line.fill.background()
    add_text(sl, Inches(1.0), y + Inches(0.05), Inches(2.2), Inches(0.25), label, 11, ACCENT, True)
    add_text(sl, Inches(3.3), y + Inches(0.05), Inches(4), Inches(0.25), desc, 11, LIGHT)
    add_text(sl, Inches(7.5), y + Inches(0.05), Inches(3.5), Inches(0.25), example, 10, GRAY)
    add_text(sl, Inches(11.2), y + Inches(0.05), Inches(1), Inches(0.25), tipo, 9, AMBER, True)
    # Actions info
    add_text(sl, Inches(1.0), y + Inches(0.32), Inches(10), Inches(0.25),
             "Todos os campos são obrigatórios", 9, GRAY)

add_info_card(sl, Inches(0.8), Inches(6.6), Inches(5.5), Inches(0.7),
              "EDITAR QUESTÃO", ["Clique no ícone ✏️ ao lado da questão existente"])
add_info_card(sl, Inches(6.8), Inches(6.6), Inches(5.5), Inches(0.7),
              "EXCLUIR QUESTÃO", ["Clique no ícone 🗑️ e confirme a exclusão"])

add_slide_number(sl, slide_num, TOTAL_SLIDES)

# ╔═══════════════════════════════════════════════════════════════╗
# ║  SLIDE 14 — SINCRONIZAÇÃO ABS                                ║
# ╚═══════════════════════════════════════════════════════════════╝
slide_num += 1
print(f"  [{slide_num:02d}] Sincronização ABS")
sl = new_slide()
add_section_header(sl, "🔄", "SINCRONIZAÇÃO DA BASE ABS", "Sobre a atualização automática da base de colaboradores")

add_text(sl, Inches(0.8), Inches(1.5), Inches(11), Inches(0.7),
         "O QUE É A SINCRONIZAÇÃO ABS?", 18, WHITE, True)

add_text(sl, Inches(0.8), Inches(2.2), Inches(11), Inches(1.2),
         "A sincronização ABS é o processo automático que atualiza a base de colaboradores da plataforma "
         "com os dados vindos do sistema ABS (Absence Management). Isso garante que a lista de "
         "funcionários esteja sempre atualizada com admissões, desligamentos e transferências.",
         13, LIGHT)

# O que faz
add_text(sl, Inches(0.8), Inches(3.5), Inches(5.5), Inches(0.4), "✅  O QUE A SINCRONIZAÇÃO FAZ:", 14, GREEN, True)
sync_does = [
    "Atualiza dados de colaboradores existentes",
    "Insere novos colaboradores contratados",
    "Remove colaboradores desligados da base ativa",
    "Promove colaboradores de Onboarding para Base Ativa",
    "Atualiza: OPSID, nome, setor, turno, líder, cargo, SOC",
]
add_bullet_list(sl, Inches(0.8), Inches(3.9), Inches(5.5), Inches(2.5), sync_does, 11, LIGHT)

# O que o PTS faz
add_text(sl, Inches(7), Inches(3.5), Inches(5.5), Inches(0.4), "👤  O QUE O PTS (ADMIN) FAZ:", 14, ACCENT, True)
pts_does = [
    "Configura o horário de sincronização automática",
    "Em Configurações → 'Sincronização Automática'",
    "Define o 'Horário de Pico (Update)'",
    "A sincronização ocorre no 1º acesso admin após esse horário",
    "Pode acionar a sincronização manual via botão 'Sincronizar Sheets'",
]
add_bullet_list(sl, Inches(7), Inches(3.9), Inches(5.5), Inches(2.5), pts_does, 11, LIGHT)

add_info_card(sl, Inches(0.8), Inches(6.2), Inches(11.5), Inches(1.0),
              "⚠️ IMPORTANTE — A SINCRONIZAÇÃO TÉCNICA É FEITA PELO DESENVOLVEDOR", [
                  "O DEV é responsável por configurar o pipeline de dados da ABS → Google Sheets → Supabase",
                  "Os PTS de cada SOC apenas configuram o horário e acionam a sincronização dentro da plataforma",
              ])

add_slide_number(sl, slide_num, TOTAL_SLIDES)

# ╔═══════════════════════════════════════════════════════════════╗
# ║  SLIDE 15 — MATERIAIS: VISÃO GERAL                          ║
# ╚═══════════════════════════════════════════════════════════════╝
slide_num += 1
print(f"  [{slide_num:02d}] Materiais - Visão Geral")
sl = new_slide()
add_section_header(sl, "📁", "MATERIAIS — VISÃO GERAL", "Repositório de pastas e links do Google Drive organizados por tema")

add_screenshot_with_border(sl, "materials", Inches(0.5), Inches(1.4), Inches(7.8), Inches(5.5))

# Anotações
add_callout_rect(sl, Inches(6.7), Inches(1.5), Inches(1.5), Inches(0.6), YELLOW_HL, Pt(3))
add_text(sl, Inches(8.5), Inches(1.5), Inches(4.5), Inches(0.3), "1️⃣ '+ Link do Google' — Adicionar material", 11, YELLOW_HL, True)

add_callout_rect(sl, Inches(6.0), Inches(1.5), Inches(0.7), Inches(0.6), GREEN, Pt(3))
add_text(sl, Inches(8.5), Inches(1.9), Inches(4.5), Inches(0.3), "2️⃣ '+ Pasta' — Criar nova pasta", 11, GREEN, True)

add_callout_rect(sl, Inches(0.8), Inches(2.1), Inches(2.5), Inches(0.5), BLUE, Pt(3))
add_text(sl, Inches(8.5), Inches(2.3), Inches(4.5), Inches(0.3), "3️⃣ Campo de busca", 11, BLUE, True)

add_callout_rect(sl, Inches(3.5), Inches(2.1), Inches(1.0), Inches(0.5), ACCENT, Pt(3))
add_text(sl, Inches(8.5), Inches(2.7), Inches(4.5), Inches(0.3), "4️⃣ Breadcrumb (navegação por pasta)", 11, ACCENT, True)

# Funcionalidades
add_text(sl, Inches(8.5), Inches(3.5), Inches(4.5), Inches(0.4), "FUNCIONALIDADES:", 14, ACCENT, True)
features = [
    "Criar pastas e subpastas (hierarquia)",
    "Adicionar links do Google Drive",
    "Busca global por nome de material",
    "Navegação por breadcrumb",
    "Renomear e excluir pastas (admin)",
    "Gerar QR Code para assinatura",
    "Excluir materiais (admin)",
]
add_bullet_list(sl, Inches(8.5), Inches(3.9), Inches(4.5), Inches(3), features, 11, LIGHT)

add_slide_number(sl, slide_num, TOTAL_SLIDES)

# ╔═══════════════════════════════════════════════════════════════╗
# ║  SLIDE 16 — MATERIAIS: CRIAR PASTA E LINK                   ║
# ╚═══════════════════════════════════════════════════════════════╝
slide_num += 1
print(f"  [{slide_num:02d}] Materiais - Criar Pasta e Link")
sl = new_slide()
add_section_header(sl, "➕", "MATERIAIS — CRIAR PASTA E ADICIONAR LINK", "Passo a passo para organizar e subir conteúdo de treinamento")

# Coluna 1 - Criar Pasta
add_text(sl, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.4), "📁  CRIAR NOVA PASTA", 16, WHITE, True)
steps_pasta = [
    "Navegue até o diretório desejado (ou Raiz)",
    "Clique no botão '+ Pasta' no canto superior direito",
    "Um formulário inline aparecerá abaixo do header",
    "Digite o nome da pasta (Ex: 'Treinamentos 2026')",
    "Clique em 'Criar' ou pressione Enter",
    "A pasta aparecerá na grade com ícone de pasta laranja",
]
for i, step in enumerate(steps_pasta):
    y = Inches(2.0 + i * 0.55)
    add_step_badge(sl, Inches(0.8), y, i + 1, Inches(0.35))
    add_text(sl, Inches(1.3), y + Inches(0.02), Inches(5), Inches(0.35), step, 11, LIGHT)

# Coluna 2 - Adicionar Link
add_text(sl, Inches(7), Inches(1.5), Inches(5.5), Inches(0.4), "🔗  ADICIONAR LINK DO GOOGLE DRIVE", 16, WHITE, True)
steps_link = [
    "Navegue até a pasta onde quer adicionar o material",
    "Clique no botão '+ Link do Google' (botão laranja)",
    "Um modal será aberto com o formulário",
    "Preencha: Nome do Treinamento (Ex: 'NR-12')",
    "Preencha: Link do Google Drive (cole o link)",
    "Clique em 'Adicionar Link'",
    "O material aparecerá na pasta com ícone azul",
]
for i, step in enumerate(steps_link):
    y = Inches(2.0 + i * 0.55)
    add_step_badge(sl, Inches(7), y, i + 1, Inches(0.35))
    add_text(sl, Inches(7.5), y + Inches(0.02), Inches(5), Inches(0.35), step, 11, LIGHT)

add_info_card(sl, Inches(0.8), Inches(5.9), Inches(5.5), Inches(1.2),
              "RENOMEAR PASTA", [
                  "Passe o mouse sobre a pasta → ícone ✏️",
                  "Um formulário inline aparece para editar",
                  "Pressione Enter ou clique 'Salvar'",
              ])

add_info_card(sl, Inches(7), Inches(5.9), Inches(5.5), Inches(1.2),
              "QR CODE DE ASSINATURA", [
                  "Passe o mouse sobre o material → ícone QR",
                  "Um modal com o QR Code é gerado automaticamente",
                  "Colaboradores escaneiam para registrar assinatura",
              ])

add_slide_number(sl, slide_num, TOTAL_SLIDES)

# ╔═══════════════════════════════════════════════════════════════╗
# ║  SLIDE 17 — COLABORADORES: VISÃO GERAL                      ║
# ╚═══════════════════════════════════════════════════════════════╝
slide_num += 1
print(f"  [{slide_num:02d}] Colaboradores - Visão Geral")
sl = new_slide()
add_section_header(sl, "👥", "COLABORADORES — VISÃO GERAL", "Base unificada de funcionários com filtros, importação e badges de onboarding")

add_screenshot_with_border(sl, "collaborators", Inches(0.5), Inches(1.4), Inches(7.8), Inches(5.5))

# Anotações
add_callout_rect(sl, Inches(2.5), Inches(1.5), Inches(2.5), Inches(0.5), YELLOW_HL, Pt(3))
add_text(sl, Inches(8.5), Inches(1.5), Inches(4.5), Inches(0.3), "1️⃣ ABAS: Base Ativa / Em Onboarding", 11, YELLOW_HL, True)

add_callout_rect(sl, Inches(5.3), Inches(1.5), Inches(3.0), Inches(0.5), GREEN, Pt(3))
add_text(sl, Inches(8.5), Inches(1.9), Inches(4.5), Inches(0.3), "2️⃣ Botões: Modelo / Importar / Sync / Novo", 11, GREEN, True)

add_callout_rect(sl, Inches(2.5), Inches(3.0), Inches(6.0), Inches(0.6), BLUE, Pt(3))
add_text(sl, Inches(8.5), Inches(2.3), Inches(4.5), Inches(0.3), "3️⃣ Filtros: Busca, SOC, Líder, Status", 11, BLUE, True)

add_callout_rect(sl, Inches(2.5), Inches(2.2), Inches(6.0), Inches(0.6), ACCENT, Pt(3))
add_text(sl, Inches(8.5), Inches(2.7), Inches(4.5), Inches(0.3), "4️⃣ Cards KPI: Total, Líderes, Treinados, %", 11, ACCENT, True)

# Features
add_text(sl, Inches(8.5), Inches(3.5), Inches(4.5), Inches(0.4), "FUNCIONALIDADES ADMIN:", 14, ACCENT, True)
feats = [
    "Adicionar colaborador manualmente",
    "Editar dados de qualquer colaborador",
    "Excluir individual ou em massa (bulk)",
    "Importar via CSV (planilha)",
    "Baixar modelo CSV de exemplo",
    "Sincronizar com Google Sheets (base ABS)",
    "Filtrar por: SOC, Líder, Status, Módulo",
    "Ver badges de onboarding (6 módulos)",
]
add_bullet_list(sl, Inches(8.5), Inches(3.9), Inches(4.5), Inches(3.5), feats, 10, LIGHT)

add_slide_number(sl, slide_num, TOTAL_SLIDES)

# ╔═══════════════════════════════════════════════════════════════╗
# ║  SLIDE 18 — COLABORADORES: NOVO REGISTRO                    ║
# ╚═══════════════════════════════════════════════════════════════╝
slide_num += 1
print(f"  [{slide_num:02d}] Colaboradores - Novo Registro")
sl = new_slide()
add_section_header(sl, "➕", "COLABORADORES — NOVO REGISTRO", "Como adicionar um colaborador manualmente")

add_text(sl, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
         "PASSO A PASSO — ADICIONAR COLABORADOR:", 16, WHITE, True)

steps = [
    "Clique no botão '+ NOVO REGISTRO' (botão vermelho no canto superior direito)",
    "Um modal será aberto com os campos adequados à aba selecionada",
    "Se estiver na aba 'Base Ativa': preencha OPSID, Gênero, Nome, BPO, Turno, Setor, Líder, Cargo, Atividade, SOC",
    "Se estiver na aba 'Em Onboarding': preencha Nome, Cargo, SOC, BPO, Data de Admissão",
    "Todos os campos de texto são automaticamente convertidos para MAIÚSCULO",
    "Clique em 'SALVAR' para cadastrar o colaborador",
]

for i, step in enumerate(steps):
    y = Inches(2.1 + i * 0.65)
    add_step_badge(sl, Inches(0.8), y, i + 1, Inches(0.4))
    add_text(sl, Inches(1.5), y + Inches(0.02), Inches(11), Inches(0.5), step, 12, LIGHT)

# Tabela de campos
add_text(sl, Inches(0.8), Inches(6.1), Inches(11), Inches(0.4), "VALIDAÇÕES APLICADAS:", 14, ACCENT, True)
validations = [
    "OPSID → Alfanumérico  |  Gênero → Só letras  |  SOC → Max 3 caracteres",
    "Turno → Max 3 chars, começa com 'T' (T1, T2)  |  Nome, Setor, Líder, Cargo → Só letras",
    "Data de Admissão → Formato dd/mm/aaaa (apenas Onboarding)",
]
for i, v in enumerate(validations):
    add_text(sl, Inches(0.8), Inches(6.5 + i * 0.3), Inches(12), Inches(0.25), f"  •  {v}", 10, GRAY)

add_slide_number(sl, slide_num, TOTAL_SLIDES)

# ╔═══════════════════════════════════════════════════════════════╗
# ║  SLIDE 19 — COLABORADORES: IMPORTAR CSV                     ║
# ╚═══════════════════════════════════════════════════════════════╝
slide_num += 1
print(f"  [{slide_num:02d}] Colaboradores - Importar CSV")
sl = new_slide()
add_section_header(sl, "📤", "COLABORADORES — IMPORTAR CSV", "Como importar colaboradores em massa via planilha")

# Coluna 1 - Passos
add_text(sl, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.4), "PASSO A PASSO:", 16, WHITE, True)
steps = [
    "Clique em 'MODELO' para baixar o CSV de exemplo",
    "Abra o CSV no Excel ou Google Sheets",
    "Preencha os dados dos colaboradores seguindo as colunas",
    "Salve o arquivo como CSV (UTF-8)",
    "Clique em 'IMPORTAR' e selecione o arquivo CSV",
    "O sistema processará em lotes de 50 registros",
    "Uma notificação mostrará o resultado da importação",
]
for i, step in enumerate(steps):
    y = Inches(2.0 + i * 0.55)
    add_step_badge(sl, Inches(0.8), y, i + 1, Inches(0.35))
    add_text(sl, Inches(1.4), y + Inches(0.02), Inches(5), Inches(0.35), step, 11, LIGHT)

# Coluna 2 - Colunas do CSV
add_text(sl, Inches(7), Inches(1.5), Inches(5.5), Inches(0.4), "COLUNAS DO CSV — BASE ATIVA:", 14, ACCENT, True)
cols_ativa = ["OPSID", "Gênero", "Colaborador", "Turno", "Setor", "Líder", "Cargo", "SOC"]
for i, col in enumerate(cols_ativa):
    y = Inches(2.0 + i * 0.4)
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), y, Inches(5.5), Inches(0.32))
    s.fill.solid(); s.fill.fore_color.rgb = DARK2; s.line.fill.background()
    add_text(sl, Inches(7.2), y + Inches(0.03), Inches(5), Inches(0.25), f"  {i+1}.  {col}", 11, LIGHT)

add_text(sl, Inches(7), Inches(5.3), Inches(5.5), Inches(0.4), "COLUNAS DO CSV — ONBOARDING:", 14, GREEN, True)
cols_onb = ["Gênero", "Colaborador", "Data de Admissão", "BPO", "Cargo", "SOC"]
for i, col in enumerate(cols_onb):
    y = Inches(5.7 + i * 0.32)
    add_text(sl, Inches(7.2), y, Inches(5), Inches(0.25), f"  {i+1}.  {col}", 11, LIGHT)

add_slide_number(sl, slide_num, TOTAL_SLIDES)

# ╔═══════════════════════════════════════════════════════════════╗
# ║  SLIDE 20 — COLABORADORES: SYNC SHEETS                      ║
# ╚═══════════════════════════════════════════════════════════════╝
slide_num += 1
print(f"  [{slide_num:02d}] Colaboradores - Sincronizar Sheets")
sl = new_slide()
add_section_header(sl, "🔄", "COLABORADORES — SINCRONIZAR SHEETS", "Atualização da base via Google Sheets (dados ABS)")

steps = [
    "Na página de Colaboradores, clique no botão '🔄 SINCRONIZAR SHEETS'",
    "O sistema buscará os dados da planilha Google Sheets publicada",
    "Matching por OPSID (prioritário) e por nome (secundário)",
    "Colaboradores novos são INSERIDOS automaticamente",
    "Colaboradores existentes são ATUALIZADOS com dados frescos",
    "Colaboradores desligados (não encontrados na planilha) são REMOVIDOS da base ativa",
    "Colaboradores em onboarding NÃO são removidos pela sincronização",
    "Ao final, uma notificação mostra: inseridos / atualizados / removidos",
]

for i, step in enumerate(steps):
    y = Inches(1.5 + i * 0.65)
    add_step_badge(sl, Inches(0.8), y, i + 1, Inches(0.4))
    add_text(sl, Inches(1.5), y + Inches(0.02), Inches(11), Inches(0.5), step, 12, LIGHT)

add_info_card(sl, Inches(0.8), Inches(6.8), Inches(5.5), Inches(0.5),
              "SYNC AUTOMÁTICO", ["Roda 1x por dia no horário configurado"])
add_info_card(sl, Inches(6.8), Inches(6.8), Inches(5.5), Inches(0.5),
              "EXCLUSIVO ADMIN", ["Apenas perfis ADMIN veem o botão"])

add_slide_number(sl, slide_num, TOTAL_SLIDES)

# ╔═══════════════════════════════════════════════════════════════╗
# ║  SLIDE 21 — COLABORADORES: BADGES ONBOARDING                ║
# ╚═══════════════════════════════════════════════════════════════╝
slide_num += 1
print(f"  [{slide_num:02d}] Colaboradores - Badges Onboarding")
sl = new_slide()
add_section_header(sl, "🏅", "BADGES DE ONBOARDING", "Sistema visual de módulos concluídos na aba Em Onboarding")

add_text(sl, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
         "Na aba 'Em Onboarding', cada colaborador exibe 6 badges de módulo:", 14, WHITE, True)

badges = [
    ("H", "Onboarding HSE", "Vermelho", "🔴", "Segurança e Saúde no Trabalho"),
    ("M", "Onboarding Meio Ambiente", "Verde", "🟢", "Normas ambientais e reciclagem"),
    ("S", "Onboarding Security", "Azul", "🔵", "Segurança patrimonial"),
    ("Q", "Onboarding Qualidade", "Âmbar", "🟡", "Controle de qualidade operacional"),
    ("R", "Onboarding People", "Roxo", "🟣", "RH, integração e cultura"),
    ("P", "Onboarding PTS", "Ciano", "🔵", "Treinamento operacional SPX"),
]

for i, (letter, name, color_name, icon, desc) in enumerate(badges):
    y = Inches(2.2 + i * 0.72)
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(11.5), Inches(0.62))
    s.fill.solid(); s.fill.fore_color.rgb = DARK2; s.line.fill.background()
    # Badge circle
    circle = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), y + Inches(0.1), Inches(0.4), Inches(0.4))
    circle.fill.solid(); circle.fill.fore_color.rgb = ACCENT; circle.line.fill.background()
    add_text(sl, Inches(1.08), y + Inches(0.12), Inches(0.25), Inches(0.3), letter, 14, WHITE, True, PP_ALIGN.CENTER)
    add_text(sl, Inches(1.6), y + Inches(0.05), Inches(0.4), Inches(0.3), icon, 14, WHITE)
    add_text(sl, Inches(2.1), y + Inches(0.05), Inches(3.5), Inches(0.3), name, 12, WHITE, True)
    add_text(sl, Inches(5.8), y + Inches(0.05), Inches(1.5), Inches(0.3), f"Cor: {color_name}", 10, GRAY)
    add_text(sl, Inches(7.5), y + Inches(0.05), Inches(4.5), Inches(0.3), desc, 11, LIGHT)

add_info_card(sl, Inches(0.8), Inches(6.6), Inches(11.5), Inches(0.7),
              "COMO FUNCIONA", [
                  "✅ Concluído = badge colorido  |  ❌ Pendente = badge cinza  |  "
                  "Conclusão vem do registro em 'trainings_completed'"])

add_slide_number(sl, slide_num, TOTAL_SLIDES)

# ╔═══════════════════════════════════════════════════════════════╗
# ║  SLIDE 22 — TREINAMENTOS EAD: VISÃO GERAL                   ║
# ╚═══════════════════════════════════════════════════════════════╝
slide_num += 1
print(f"  [{slide_num:02d}] Treinamentos EAD - Visão Geral")
sl = new_slide()
add_section_header(sl, "🎓", "CAPACITAÇÃO EAD — VISÃO GERAL", "Plataforma de treinamento com vídeo, prova e assinatura digital")

add_screenshot_with_border(sl, "trainings", Inches(0.5), Inches(1.4), Inches(7.8), Inches(5.5))

# Anotações
add_callout_rect(sl, Inches(5.5), Inches(1.5), Inches(1.3), Inches(0.5), GREEN, Pt(3))
add_text(sl, Inches(8.5), Inches(1.5), Inches(4.5), Inches(0.3), "1️⃣ 'NOVA PASTA' — Criar pasta de curso", 11, GREEN, True)

add_callout_rect(sl, Inches(6.8), Inches(1.5), Inches(1.5), Inches(0.5), YELLOW_HL, Pt(3))
add_text(sl, Inches(8.5), Inches(1.9), Inches(4.5), Inches(0.3), "2️⃣ 'NOVO MATERIAL' — Adicionar vídeo", 11, YELLOW_HL, True)

add_callout_rect(sl, Inches(2.4), Inches(1.9), Inches(0.8), Inches(0.3), BLUE, Pt(3))
add_text(sl, Inches(8.5), Inches(2.3), Inches(4.5), Inches(0.3), "3️⃣ Breadcrumb 'RAIZ' — Navegação", 11, BLUE, True)

# Funcionalidades
add_text(sl, Inches(8.5), Inches(3.0), Inches(4.5), Inches(0.4), "FUNCIONALIDADES ADMIN:", 14, ACCENT, True)
feats = [
    "Criar pastas para organizar treinamentos",
    "Adicionar materiais com link de vídeo (Google Drive)",
    "Excluir pastas e materiais",
    "Realizar treinamentos (vídeo + prova + assinatura)",
    "Navegação hierárquica via breadcrumb",
]
add_bullet_list(sl, Inches(8.5), Inches(3.4), Inches(4.5), Inches(2), feats, 11, LIGHT)

add_text(sl, Inches(8.5), Inches(5.3), Inches(4.5), Inches(0.4), "FLUXO DE CERTIFICAÇÃO:", 14, ACCENT, True)
flow = [
    "Fase 1: Assistir ao vídeo (60s mínimo)",
    "Fase 2: Fazer a prova (90% para aprovar)",
    "Fase 3: Assinatura digital no canvas",
    "Resultado: Certificação registrada ✅",
]
add_numbered_list(sl, Inches(8.5), Inches(5.7), Inches(4.5), Inches(1.5), flow, 11, LIGHT)

add_slide_number(sl, slide_num, TOTAL_SLIDES)

# ╔═══════════════════════════════════════════════════════════════╗
# ║  SLIDE 23 — TREINAMENTOS: CRIAR MATERIAL                    ║
# ╚═══════════════════════════════════════════════════════════════╝
slide_num += 1
print(f"  [{slide_num:02d}] Treinamentos - Criar Material")
sl = new_slide()
add_section_header(sl, "➕", "CAPACITAÇÃO EAD — CRIAR MATERIAL", "Como adicionar um novo treinamento com vídeo do Google Drive")

add_text(sl, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
         "PASSO A PASSO — ADICIONAR TREINAMENTO:", 16, WHITE, True)

steps = [
    ("Navegue até a pasta desejada (ou crie uma nova)",
     "Use os botões 'NOVA PASTA' para organizar por tema"),
    ("Clique no botão '+ NOVO MATERIAL' (botão vermelho)",
     "Um modal será aberto com o formulário de criação"),
    ("Preencha o 'Título do Conteúdo'",
     "Ex: 'Onboarding SPX', 'NR-12 Segurança', 'Processamento Básico'"),
    ("Cole o 'Link do Vídeo (Google Drive)'",
     "Copie o link de compartilhamento do Google Drive do vídeo"),
    ("Clique em 'Salvar Treinamento'",
     "O material aparecerá com ícone de play verde na pasta"),
    ("Configure as questões da prova em Configurações",
     "Vá em Configurações → Avaliações e Provas → selecione o treinamento → adicione questões"),
]

for i, (step, detail) in enumerate(steps):
    y = Inches(2.1 + i * 0.85)
    add_step_badge(sl, Inches(0.8), y, i + 1, Inches(0.42))
    add_text(sl, Inches(1.5), y + Inches(0.02), Inches(5.5), Inches(0.35), step, 13, WHITE, True)
    add_text(sl, Inches(1.5), y + Inches(0.38), Inches(5.5), Inches(0.3), detail, 11, GRAY)

add_info_card(sl, Inches(7.5), Inches(1.5), Inches(5), Inches(2),
              "FORMATO DO LINK DO VÍDEO", [
                  "Use links de COMPARTILHAMENTO do Google Drive",
                  "Ex: https://drive.google.com/file/d/XXXXX/view",
                  "O sistema converte automaticamente para embed",
                  "Se o link não iniciar com http, 'https://' é adicionado",
              ])

add_info_card(sl, Inches(7.5), Inches(3.8), Inches(5), Inches(1.8),
              "IMPORTANTE", [
                  "O treinamento SÓ terá prova se você configurar",
                  "questões em Configurações → Avaliações",
                  "Sem questões: aparece 'CONTEÚDO SEM AVALIAÇÃO'",
                  "Com questões: colaborador precisa 90% para aprovar",
              ])

add_info_card(sl, Inches(7.5), Inches(5.9), Inches(5), Inches(1.2),
              "EXCLUIR MATERIAL", [
                  "Passe o mouse sobre o card do material",
                  "Clique no ícone 🗑️ que aparece",
                  "Confirme a exclusão no diálogo",
              ])

add_slide_number(sl, slide_num, TOTAL_SLIDES)

# ╔═══════════════════════════════════════════════════════════════╗
# ║  SLIDE 24 — AGENDA: VISÃO GERAL                             ║
# ╚═══════════════════════════════════════════════════════════════╝
slide_num += 1
print(f"  [{slide_num:02d}] Agenda - Visão Geral")
sl = new_slide()
add_section_header(sl, "📅", "AGENDA DE TREINAMENTOS — VISÃO GERAL", "Calendário semanal para gerenciar treinamentos presenciais")

add_screenshot_with_border(sl, "schedule", Inches(0.5), Inches(1.4), Inches(7.8), Inches(5.5))

# Anotações
add_callout_rect(sl, Inches(4.8), Inches(1.5), Inches(1.2), Inches(0.5), YELLOW_HL, Pt(3))
add_text(sl, Inches(8.5), Inches(1.5), Inches(4.5), Inches(0.3), "1️⃣ Filtro por SOC", 11, YELLOW_HL, True)

add_callout_rect(sl, Inches(5.3), Inches(1.5), Inches(1.5), Inches(0.5), ACCENT, Pt(3))
add_text(sl, Inches(8.5), Inches(1.9), Inches(4.5), Inches(0.3), "2️⃣ Abas: Calendário / Solicitações / Histórico", 11, ACCENT, True)

add_callout_rect(sl, Inches(7.0), Inches(1.5), Inches(1.0), Inches(0.5), GREEN, Pt(3))
add_text(sl, Inches(8.5), Inches(2.3), Inches(4.5), Inches(0.3), "3️⃣ Navegação semanal (< Hoje >)", 11, GREEN, True)

add_callout_rect(sl, Inches(7.5), Inches(1.5), Inches(1.0), Inches(0.5), BLUE, Pt(3))
add_text(sl, Inches(8.5), Inches(2.7), Inches(4.5), Inches(0.3), "4️⃣ Botão 'Configurar' (ADMIN only)", 11, BLUE, True)

# Features
add_text(sl, Inches(8.5), Inches(3.4), Inches(4.5), Inches(0.4), "FUNCIONALIDADES:", 14, ACCENT, True)
feats = [
    "Calendário semanal visual (06h-22h)",
    "Filtro por SOC",
    "3 abas: Calendário, Solicitações, Histórico",
    "Configurar slots (admin) — criar/editar/excluir",
    "Inscrever colaboradores nos treinamentos",
    "Regra D+2 (inscrição com 2 dias de antecedência)",
    "Integração com Google Calendar",
    "Auditoria completa (quem inscreveu/removeu)",
    "Solicitações de líderes com aprovação PTS",
]
add_bullet_list(sl, Inches(8.5), Inches(3.8), Inches(4.5), Inches(3.5), feats, 10, LIGHT)

add_slide_number(sl, slide_num, TOTAL_SLIDES)

# ╔═══════════════════════════════════════════════════════════════╗
# ║  SLIDE 25 — AGENDA: CRIAR SLOT                              ║
# ╚═══════════════════════════════════════════════════════════════╝
slide_num += 1
print(f"  [{slide_num:02d}] Agenda - Criar Slot")
sl = new_slide()
add_section_header(sl, "⚡", "AGENDA — CRIAR SLOT DE TREINAMENTO", "Como configurar um novo horário de treinamento no calendário")

add_text(sl, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
         "PASSO A PASSO — CRIAR SLOT:", 16, WHITE, True)

steps = [
    ("Clique no botão '⚙ Configurar' no canto superior direito",
     "O painel administrativo será aberto abaixo do header"),
    ("Clique em '+ Novo Slot'",
     "Um formulário será exibido com os campos do agendamento"),
    ("Preencha o TÍTULO do treinamento",
     "Ex: 'Treinamento NR-12 — Recebimento'"),
    ("Selecione o TIPO de treinamento",
     "Opções: RECEBIMENTO, PROCESSAMENTO, EXPEDIÇÃO, etc."),
    ("Defina se é RECORRENTE ou DATA ESPECÍFICA",
     "Recorrente: selecione dia da semana / Específica: selecione a data"),
    ("Defina HORÁRIO de início e fim",
     "Ex: 10:00 às 11:00"),
    ("Preencha INSTRUTOR (nome e e-mail)",
     "O e-mail é usado para envio do convite Google Calendar"),
    ("Selecione o LOCAL e a SOC",
     "Ex: Local = 'Sala de Treinamento A', SOC = 'SP6'"),
    ("Clique em 'SALVAR'",
     "O slot aparecerá no calendário semanal no horário configurado"),
]

for i, (step, detail) in enumerate(steps):
    y = Inches(2.0 + i * 0.58)
    add_step_badge(sl, Inches(0.8), y, i + 1, Inches(0.35))
    add_text(sl, Inches(1.3), y + Inches(0.02), Inches(5.5), Inches(0.3), step, 11, WHITE, True)
    add_text(sl, Inches(7.0), y + Inches(0.02), Inches(5.5), Inches(0.3), detail, 10, GRAY)

add_slide_number(sl, slide_num, TOTAL_SLIDES)

# ╔═══════════════════════════════════════════════════════════════╗
# ║  SLIDE 26 — AGENDA: INSCREVER COLABORADORES                 ║
# ╚═══════════════════════════════════════════════════════════════╝
slide_num += 1
print(f"  [{slide_num:02d}] Agenda - Inscrever Colaboradores")
sl = new_slide()
add_section_header(sl, "📋", "AGENDA — INSCREVER COLABORADORES", "Como inscrever e remover colaboradores de um slot de treinamento")

steps = [
    "No calendário, clique sobre o SLOT desejado (bloco colorido no dia/hora)",
    "Um painel lateral será aberto mostrando os detalhes do treinamento",
    "Veja os colaboradores JÁ INSCRITOS na lista do painel",
    "Use o campo de busca para localizar o colaborador desejado",
    "Use os filtros de: Área (setor), Cargo, Status (treinado/pendente)",
    "Clique no nome do colaborador para INSCREVÊ-LO",
    "Para REMOVER: clique no ícone 🗑️ ao lado do nome inscrito",
    "A inscrição gera automaticamente um registro de auditoria",
]

for i, step in enumerate(steps):
    y = Inches(1.5 + i * 0.6)
    add_step_badge(sl, Inches(0.8), y, i + 1, Inches(0.4))
    add_text(sl, Inches(1.5), y + Inches(0.02), Inches(7), Inches(0.45), step, 12, LIGHT)

add_info_card(sl, Inches(9), Inches(1.5), Inches(3.8), Inches(1.5),
              "REGRA D+2", [
                  "Só é possível inscrever em datas",
                  "com no mínimo 2 dias de antecedência",
                  "Isso garante tempo para organização",
              ])

add_info_card(sl, Inches(9), Inches(3.3), Inches(3.8), Inches(1.5),
              "GOOGLE CALENDAR", [
                  "Ao inscrever, um convite é enviado",
                  "automaticamente ao instrutor via",
                  "Google Calendar com todos os detalhes",
              ])

add_info_card(sl, Inches(9), Inches(5.1), Inches(3.8), Inches(1.5),
              "FILTROS DISPONÍVEIS", [
                  "Área — Setor operacional",
                  "Cargo — Função do colaborador",
                  "Status — Certificado ou Pendente",
              ])

add_slide_number(sl, slide_num, TOTAL_SLIDES)

# ╔═══════════════════════════════════════════════════════════════╗
# ║  SLIDE 27 — AGENDA: SOLICITAÇÕES (VISÃO)                    ║
# ╚═══════════════════════════════════════════════════════════════╝
slide_num += 1
print(f"  [{slide_num:02d}] Agenda - Solicitações")
sl = new_slide()
add_section_header(sl, "📨", "AGENDA — SOLICITAÇÕES DE TREINAMENTO",
                   "Como os líderes enviam solicitações e como o ADMIN analisa")

# Screenshot da lista de solicitações
add_screenshot_with_border(sl, "pts_requests", Inches(0.5), Inches(1.4), Inches(5.5), Inches(5.5))

# Anotações
add_text(sl, Inches(6.5), Inches(1.5), Inches(6.3), Inches(0.4),
         "FLUXO DE SOLICITAÇÃO:", 16, ACCENT, True)

flow_steps = [
    ("Líder cria solicitação", "Na aba 'Solicitações', clica em '+ Requisição'"),
    ("Preenche o formulário", "Tipo de treinamento, data (D+2), horário, colaboradores"),
    ("Solicitação fica 'PENDENTE'", "Aparece para o ADMIN com botão 'Analisar'"),
    ("ADMIN analisa a solicitação", "Verifica detalhes, pode excluir colaboradores"),
    ("ADMIN aprova ou recusa", "Aprovar: define instrutor / Recusar: informa motivo"),
    ("Se APROVADO", "Slot é criado no calendário + inscrições automáticas"),
]

for i, (title, desc) in enumerate(flow_steps):
    y = Inches(2.0 + i * 0.72)
    add_step_badge(sl, Inches(6.5), y, i + 1, Inches(0.38))
    add_text(sl, Inches(7.1), y + Inches(0.0), Inches(5.5), Inches(0.3), title, 12, WHITE, True)
    add_text(sl, Inches(7.1), y + Inches(0.32), Inches(5.5), Inches(0.3), desc, 10, GRAY)

# Status badges
add_text(sl, Inches(6.5), Inches(6.5), Inches(6), Inches(0.3), "STATUS POSSÍVEIS:", 12, ACCENT, True)
statuses = [
    ("⏳ PENDENTE", "Aguardando análise do PTS/Admin", AMBER),
    ("✅ APROVADA", "Treinamento confirmado e publicado", GREEN),
    ("❌ RECUSADA", "Solicitação negada com motivo", RED),
]
for i, (badge, desc, color) in enumerate(statuses):
    add_text(sl, Inches(6.5 + i * 2.2), Inches(6.8), Inches(2), Inches(0.3), badge, 10, color, True)
    add_text(sl, Inches(6.5 + i * 2.2), Inches(7.1), Inches(2), Inches(0.3), desc, 8, GRAY)

add_slide_number(sl, slide_num, TOTAL_SLIDES)

# ╔═══════════════════════════════════════════════════════════════╗
# ║  SLIDE 28 — AGENDA: APROVAR SOLICITAÇÃO                     ║
# ╚═══════════════════════════════════════════════════════════════╝
slide_num += 1
print(f"  [{slide_num:02d}] Agenda - Aprovar Solicitação")
sl = new_slide()
add_section_header(sl, "✅", "AGENDA — ANALISAR E APROVAR SOLICITAÇÃO",
                   "Passo a passo para aprovar ou recusar uma solicitação de treinamento")

# Screenshot do modal de aprovação
add_screenshot_with_border(sl, "pts_approval", Inches(0.5), Inches(1.4), Inches(5.5), Inches(5.5))

# Anotações
add_callout_rect(sl, Inches(1.3), Inches(2.5), Inches(3.8), Inches(1.2), YELLOW_HL, Pt(3))
add_text(sl, Inches(6.5), Inches(1.5), Inches(6.3), Inches(0.4),
         "ANALISANDO UMA SOLICITAÇÃO:", 16, ACCENT, True)

analysis_steps = [
    "Na aba 'Solicitações', clique em 'Analisar' na solicitação PENDENTE",
    "O modal 'Analisar Solicitação' é aberto com todos os detalhes",
    "Verifique: Tipo de Treinamento, Data, Horário, Local, Solicitante",
    "Veja a lista de colaboradores incluídos na solicitação",
    "Se necessário, EXCLUA colaboradores individuais (ícone ⊖)",
    "Para APROVAR: clique em '✓ Aprovar'",
    "   → Será solicitado o NOME e E-MAIL do instrutor",
    "   → O sistema cria o slot + inscreve todos automaticamente",
    "Para RECUSAR: clique em '✗ Recusar'",
    "   → Informe o MOTIVO da recusa (obrigatório)",
]

for i, step in enumerate(analysis_steps):
    y = Inches(2.0 + i * 0.48)
    if not step.startswith("   "):
        add_text(sl, Inches(6.5), y, Inches(6.3), Inches(0.35), f"  {i+1}.  {step}", 11, LIGHT)
    else:
        add_text(sl, Inches(6.5), y, Inches(6.3), Inches(0.35), step, 11, GRAY)

add_slide_number(sl, slide_num, TOTAL_SLIDES)

# ╔═══════════════════════════════════════════════════════════════╗
# ║  SLIDE 29 — AGENDA: FORMULÁRIO DE SOLICITAÇÃO               ║
# ╚═══════════════════════════════════════════════════════════════╝
slide_num += 1
print(f"  [{slide_num:02d}] Agenda - Formulário de Solicitação")
sl = new_slide()
add_section_header(sl, "📝", "AGENDA — FORMULÁRIO DE SOLICITAÇÃO",
                   "Campos do formulário que o líder preenche ao solicitar um treinamento")

add_screenshot_with_border(sl, "pts_form", Inches(0.5), Inches(1.4), Inches(5.5), Inches(5.5))

# Campos explicados
add_text(sl, Inches(6.5), Inches(1.5), Inches(6.3), Inches(0.4),
         "CAMPOS DO FORMULÁRIO:", 16, ACCENT, True)

form_fields = [
    ("TIPO DE TREINAMENTO", "Dropdown: RECEBIMENTO, PROCESSAMENTO, EXPEDIÇÃO, etc."),
    ("DATA PREVISTA", "Selecionar data (mínimo D+2 a partir de hoje)"),
    ("HORA DE INÍCIO", "Horário do início do treinamento (ex: 14:00)"),
    ("HORA DE FIM", "Calculado automaticamente (início + 1h)"),
    ("LOCAL", "Preenchido automaticamente com a SOC do solicitante"),
    ("SELECIONAR COLABORADORES", "Busca e seleciona os colaboradores participantes"),
]

for i, (label, desc) in enumerate(form_fields):
    y = Inches(2.0 + i * 0.72)
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), y, Inches(6.3), Inches(0.62))
    s.fill.solid(); s.fill.fore_color.rgb = DARK2; s.line.fill.background()
    add_text(sl, Inches(6.7), y + Inches(0.05), Inches(2.5), Inches(0.25), label, 10, ACCENT, True)
    add_text(sl, Inches(6.7), y + Inches(0.3), Inches(5.8), Inches(0.25), desc, 11, LIGHT)

add_info_card(sl, Inches(6.5), Inches(6.4), Inches(6.3), Inches(0.8),
              "⚠️ REGRA D+2", [
                  "A data selecionada deve ser no mínimo 2 dias após a data atual. Solicitações com data menor são rejeitadas."
              ])

add_slide_number(sl, slide_num, TOTAL_SLIDES)

# ╔═══════════════════════════════════════════════════════════════╗
# ║  SLIDE 30 — SOCs: GESTÃO DE UNIDADES                        ║
# ╚═══════════════════════════════════════════════════════════════╝
slide_num += 1
print(f"  [{slide_num:02d}] SOCs - Gestão de Unidades")
sl = new_slide()
add_section_header(sl, "🏢", "SOCs — GESTÃO DE UNIDADES", "Como cadastrar e gerenciar as unidades operacionais")

# Layout
add_text(sl, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.4), "COMO CADASTRAR UMA NOVA SOC:", 16, WHITE, True)

steps = [
    "Acesse o módulo 'SOCs' na sidebar",
    "Clique no botão '+ Nova SOC' (canto superior direito)",
    "No modal, preencha os campos:",
    "   → Nome da Unidade (ex: SP6, SP1, RJ1)",
    "   → Endereço Completo",
    "   → PTS Responsável (nome do PTS local)",
    "   → Site Leader (nome do gestor)",
    "   → Marque se possui processo de Sorting",
    "Clique em 'Cadastrar Unidade'",
    "A SOC aparecerá na grade de unidades",
]
for i, step in enumerate(steps):
    y = Inches(2.0 + i * 0.48)
    if not step.startswith("   "):
        add_step_badge(sl, Inches(0.8), y, i + 1 if not step.startswith("   ") else 0, Inches(0.3))
    add_text(sl, Inches(1.3) if not step.startswith("   ") else Inches(1.3), y + Inches(0.01), Inches(5), Inches(0.3),
             step, 11, LIGHT if not step.startswith("   ") else GRAY)

# Campos do formulário
add_text(sl, Inches(7), Inches(1.5), Inches(5.5), Inches(0.4), "CAMPOS DO CADASTRO:", 14, ACCENT, True)

soc_fields = [
    ("NOME DA UNIDADE", "Identificação da SOC (ex: SP6)", "Obrigatório"),
    ("ENDEREÇO", "Endereço físico da unidade", "Opcional"),
    ("PTS RESPONSÁVEL", "Nome do PTS daquela base", "Opcional"),
    ("SITE LEADER", "Gestor responsável pela unidade", "Opcional"),
    ("SORTING", "Checkbox: possui processo de sorting?", "Sim/Não"),
]
for i, (label, desc, obs) in enumerate(soc_fields):
    y = Inches(2.0 + i * 0.7)
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), y, Inches(5.5), Inches(0.58))
    s.fill.solid(); s.fill.fore_color.rgb = DARK2; s.line.fill.background()
    add_text(sl, Inches(7.2), y + Inches(0.05), Inches(2.5), Inches(0.25), label, 10, ACCENT, True)
    add_text(sl, Inches(7.2), y + Inches(0.3), Inches(3.5), Inches(0.25), desc, 10, LIGHT)
    add_text(sl, Inches(11), y + Inches(0.05), Inches(1.3), Inches(0.25), obs, 9, AMBER, True)

add_info_card(sl, Inches(7), Inches(5.6), Inches(5.5), Inches(1.5),
              "GERENCIAR SOCs", [
                  "EDITAR: passe o mouse → ícone ✏️",
                  "EXCLUIR: passe o mouse → ícone 🗑️",
                  "VER DETALHES: clique no card da SOC",
                  "⚠️ Excluir uma SOC pode afetar dados vinculados!",
              ])

add_slide_number(sl, slide_num, TOTAL_SLIDES)

# ╔═══════════════════════════════════════════════════════════════╗
# ║  SLIDE 31 — ASSINATURAS REGISTRADAS                         ║
# ╚═══════════════════════════════════════════════════════════════╝
slide_num += 1
print(f"  [{slide_num:02d}] Assinaturas Registradas")
sl = new_slide()
add_section_header(sl, "✍️", "ASSINATURAS REGISTRADAS", "Auditoria completa de todas as assinaturas digitais de treinamento")

add_screenshot_with_border(sl, "signatures", Inches(0.5), Inches(1.4), Inches(7.8), Inches(5.5))

# Colunas
add_text(sl, Inches(8.8), Inches(1.5), Inches(4.2), Inches(0.4), "INFORMAÇÕES EXIBIDAS:", 14, ACCENT, True)

cols = [
    ("COLABORADOR", "Nome e cargo do funcionário"),
    ("TREINAMENTO", "Tipo do treinamento concluído (badge colorido)"),
    ("INSTRUTOR", "Quem ministrou o treinamento"),
    ("SOC", "Unidade operacional"),
    ("DATA", "Data do registro da assinatura"),
]
for i, (col, desc) in enumerate(cols):
    y = Inches(2.0 + i * 0.6)
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), y, Inches(4.2), Inches(0.5))
    s.fill.solid(); s.fill.fore_color.rgb = DARK2; s.line.fill.background()
    add_text(sl, Inches(9.0), y + Inches(0.05), Inches(2), Inches(0.2), col, 10, WHITE, True)
    add_text(sl, Inches(9.0), y + Inches(0.25), Inches(4), Inches(0.2), desc, 10, GRAY)

add_text(sl, Inches(8.8), Inches(5.1), Inches(4.2), Inches(0.4), "FILTROS DISPONÍVEIS:", 14, ACCENT, True)
add_bullet_list(sl, Inches(8.8), Inches(5.5), Inches(4.2), Inches(1.5), [
    "Busca por colaborador, treinamento ou instrutor",
    "Filtro por SOC (todas ou uma específica)",
    "Download individual da assinatura (imagem)",
], 11, LIGHT)

add_slide_number(sl, slide_num, TOTAL_SLIDES)

# ╔═══════════════════════════════════════════════════════════════╗
# ║  SLIDE 32 — RELATÓRIOS                                      ║
# ╚═══════════════════════════════════════════════════════════════╝
slide_num += 1
print(f"  [{slide_num:02d}] Relatórios")
sl = new_slide()
add_section_header(sl, "📈", "RELATÓRIOS E MATRIZ", "Gestão de certificações por unidade e setor operacional")

add_screenshot_with_border(sl, "reports", Inches(0.5), Inches(1.4), Inches(7.8), Inches(5.5))

# Anotações
add_callout_rect(sl, Inches(2.5), Inches(1.7), Inches(6.0), Inches(1.8), YELLOW_HL, Pt(2))
add_text(sl, Inches(8.8), Inches(1.5), Inches(4.2), Inches(0.3), "1️⃣ Cards de % por setor (Geral a Tratativas)", 11, YELLOW_HL, True)

add_callout_rect(sl, Inches(4.8), Inches(1.5), Inches(1.3), Inches(0.5), GREEN, Pt(2))
add_text(sl, Inches(8.8), Inches(1.9), Inches(4.2), Inches(0.3), "2️⃣ Filtro por unidade (SOC)", 11, GREEN, True)

add_callout_rect(sl, Inches(6.3), Inches(1.5), Inches(2.0), Inches(0.5), BLUE, Pt(2))
add_text(sl, Inches(8.8), Inches(2.3), Inches(4.2), Inches(0.3), "3️⃣ Filtro por período e status", 11, BLUE, True)

# Features
add_text(sl, Inches(8.8), Inches(3.0), Inches(4.2), Inches(0.4), "TIPOS DE RELATÓRIO:", 14, ACCENT, True)
reports = [
    "Desempenho geral por setor (%)",
    "Desempenho por SOC (gráficos)",
    "Matriz de Onboarding",
    "Certificação Operacional",
    "Ranking de Instrutores",
    "Filtros por período e status",
]
add_bullet_list(sl, Inches(8.8), Inches(3.4), Inches(4.2), Inches(2.5), reports, 11, LIGHT)

add_info_card(sl, Inches(8.8), Inches(5.8), Inches(4.2), Inches(1.2),
              "COMO USAR", [
                  "Selecione a SOC no dropdown superior",
                  "Defina o período desejado (opcional)",
                  "Filtre por status: Todos/Certificados/Pendentes",
              ])

add_slide_number(sl, slide_num, TOTAL_SLIDES)

# ╔═══════════════════════════════════════════════════════════════╗
# ║  SLIDE 33 — TABELA DE PERMISSÕES POR PERFIL                 ║
# ╚═══════════════════════════════════════════════════════════════╝
slide_num += 1
print(f"  [{slide_num:02d}] Tabela de Permissões")
sl = new_slide()
add_section_header(sl, "🔒", "MAPA DE PERMISSÕES POR PERFIL", "Tabela completa de acesso por módulo e por perfil")

# Header da tabela
headers = ["MÓDULO", "ADMIN", "USUÁRIO", "LÍDER", "BPO", "PCP"]
col_widths = [Inches(3.5), Inches(1.4), Inches(1.4), Inches(1.4), Inches(1.4), Inches(1.4)]
x_start = Inches(0.8)
y_header = Inches(1.4)

for i, (hdr, w) in enumerate(zip(headers, col_widths)):
    x = x_start + sum(cw for cw in col_widths[:i])
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y_header, w, Inches(0.45))
    s.fill.solid(); s.fill.fore_color.rgb = ACCENT; s.line.fill.background()
    add_text(sl, x + Inches(0.1), y_header + Inches(0.08), w - Inches(0.1), Inches(0.3),
             hdr, 10, WHITE, True, PP_ALIGN.CENTER)

# Data rows
permissions = [
    ("Dashboard", "✅", "✅", "✅", "❌", "❌"),
    ("Materiais", "✅", "✅", "❌", "❌", "❌"),
    ("Colaboradores", "✅", "✅", "✅ (time)", "✅ (onb)", "❌"),
    ("Relatórios", "✅", "✅", "✅", "❌", "❌"),
    ("Agenda", "✅", "✅", "✅", "❌", "✅"),
    ("Assinaturas", "✅", "✅", "❌", "❌", "❌"),
    ("SOCs", "✅", "✅", "❌", "❌", "❌"),
    ("Treinamentos EAD", "✅", "✅", "✅", "❌", "❌"),
    ("Configurações", "✅", "❌", "❌", "❌", "❌"),
    ("Criar Usuários", "✅", "❌", "❌", "❌", "❌"),
    ("Criar Instrutores", "✅", "❌", "❌", "❌", "❌"),
    ("Gerenciar Provas", "✅", "❌", "❌", "❌", "❌"),
    ("Sincronizar Base", "✅", "❌", "❌", "❌", "❌"),
    ("Excluir Colaboradores", "✅", "❌", "❌", "❌", "❌"),
]

for row_i, (mod, *perms) in enumerate(permissions):
    y = y_header + Inches(0.45) + row_i * Inches(0.38)
    bg = DARK2 if row_i % 2 == 0 else DARK3
    for col_i, (val, w) in enumerate(zip([mod] + list(perms), col_widths)):
        x = x_start + sum(cw for cw in col_widths[:col_i])
        s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.38))
        s.fill.solid(); s.fill.fore_color.rgb = bg; s.line.fill.background()
        color = WHITE if col_i == 0 else (GREEN if "✅" in val else RED if "❌" in val else LIGHT)
        bold = col_i == 0
        add_text(sl, x + Inches(0.1), y + Inches(0.05), w - Inches(0.1), Inches(0.25),
                 val, 9 if col_i > 0 else 10, color, bold, PP_ALIGN.CENTER if col_i > 0 else PP_ALIGN.LEFT)

add_slide_number(sl, slide_num, TOTAL_SLIDES)

# ╔═══════════════════════════════════════════════════════════════╗
# ║  SLIDE 34 — RESUMO: CHECKLIST DO ADMIN                      ║
# ╚═══════════════════════════════════════════════════════════════╝
slide_num += 1
print(f"  [{slide_num:02d}] Checklist do Admin")
sl = new_slide()
add_section_header(sl, "📋", "CHECKLIST DO ADMIN", "O que fazer ao configurar a plataforma para uma nova SOC")

add_text(sl, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
         "Siga esta ordem para configurar a plataforma corretamente:", 14, LIGHT)

checklist = [
    ("⬜", "Cadastrar a SOC", "SOCs → + Nova SOC (nome, endereço, PTS, site leader)"),
    ("⬜", "Criar o usuário Admin/PTS da SOC", "Configurações → + Usuário (e-mail, senha, perfil ADMIN, SOC)"),
    ("⬜", "Cadastrar os instrutores", "Configurações → + Instrutor (nome, SOC)"),
    ("⬜", "Criar pastas de materiais", "Materiais → + Pasta (organizar por tema/área)"),
    ("⬜", "Adicionar links do Google Drive", "Materiais → + Link do Google (links de apresentações/documentos)"),
    ("⬜", "Criar pastas de treinamentos EAD", "Treinamentos → + Nova Pasta → + Novo Material (vídeos)"),
    ("⬜", "Configurar as provas de certificação", "Configurações → Avaliações → Adicionar questões por treinamento"),
    ("⬜", "Importar base de colaboradores", "Colaboradores → Importar CSV ou Sincronizar Sheets"),
    ("⬜", "Configurar horário de sincronização", "Configurações → Sincronização Automática → Horário"),
    ("⬜", "Configurar a agenda de treinamentos", "Agenda → Configurar → + Novo Slot (recorrente ou específico)"),
    ("⬜", "Criar usuários para líderes (se necessário)", "Configurações → + Usuário (perfil 'Líder Operacional')"),
    ("⬜", "Testar o fluxo completo", "Logar como líder, fazer treinamento EAD, verificar assinatura"),
]

for i, (icon, title, desc) in enumerate(checklist):
    y = Inches(2.0 + i * 0.43)
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), y, Inches(12), Inches(0.38))
    s.fill.solid(); s.fill.fore_color.rgb = DARK2 if i % 2 == 0 else DARK3; s.line.fill.background()
    add_text(sl, Inches(0.8), y + Inches(0.04), Inches(0.3), Inches(0.25), icon, 11, WHITE)
    add_text(sl, Inches(1.2), y + Inches(0.04), Inches(3.5), Inches(0.25), title, 11, WHITE, True)
    add_text(sl, Inches(4.8), y + Inches(0.04), Inches(7.5), Inches(0.25), desc, 10, GRAY)

add_slide_number(sl, slide_num, TOTAL_SLIDES)

# ╔═══════════════════════════════════════════════════════════════╗
# ║  SLIDE 35 — DICAS E BOAS PRÁTICAS                           ║
# ╚═══════════════════════════════════════════════════════════════╝
slide_num += 1
print(f"  [{slide_num:02d}] Dicas e Boas Práticas")
sl = new_slide()
add_section_header(sl, "💡", "DICAS E BOAS PRÁTICAS", "Recomendações para uso eficiente da plataforma")

tips = [
    ("🔐", "Senhas seguras", "Use senhas provisórias aleatórias e oriente a troca no primeiro acesso"),
    ("📁", "Organize os materiais", "Crie pastas por área (COP, HSE, People, PTS, etc.) para fácil localização"),
    ("📝", "Provas bem feitas", "Mínimo 5 questões por prova — cubra os pontos essenciais do treinamento"),
    ("🔄", "Sincronize diariamente", "Configure o horário de sincronização para antes do expediente (ex: 05:00)"),
    ("📅", "Agenda organizada", "Use slots recorrentes para treinamentos regulares e datas específicas para especiais"),
    ("👤", "Chave do Líder", "Ao criar perfil de Líder, preencha a Chave de Identificação EXATAMENTE como na planilha"),
    ("📊", "Monitore relatórios", "Verifique semanalmente os % de certificação por setor para identificar gaps"),
    ("🏢", "SOCs completas", "Preencha todos os dados da SOC (endereço, PTS, site leader) para documentação"),
    ("✅", "Teste antes de treinar", "Faça você mesmo o fluxo completo (vídeo → prova → assinatura) antes de apresentar"),
]

for i, (icon, title, desc) in enumerate(tips):
    y = Inches(1.5 + i * 0.62)
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(11.5), Inches(0.52))
    s.fill.solid(); s.fill.fore_color.rgb = DARK2; s.line.fill.background()
    add_text(sl, Inches(1.0), y + Inches(0.05), Inches(0.4), Inches(0.4), icon, 16, WHITE)
    add_text(sl, Inches(1.6), y + Inches(0.05), Inches(3), Inches(0.35), title, 12, WHITE, True)
    add_text(sl, Inches(4.8), y + Inches(0.05), Inches(7.3), Inches(0.35), desc, 11, GRAY)

add_slide_number(sl, slide_num, TOTAL_SLIDES)

# ╔═══════════════════════════════════════════════════════════════╗
# ║  SLIDE 36 — TROUBLESHOOTING                                  ║
# ╚═══════════════════════════════════════════════════════════════╝
slide_num += 1
print(f"  [{slide_num:02d}] Troubleshooting")
sl = new_slide()
add_section_header(sl, "🔧", "RESOLUÇÃO DE PROBLEMAS", "Problemas comuns e como resolvê-los")

problems = [
    ("Não consigo ver 'Configurações'", "Seu perfil não é ADMIN. Peça ao administrador para alterar seu perfil para 'admin'"),
    ("Usuário não consegue logar", "Verifique: e-mail correto? Senha provisória informada? Tente recriar o usuário"),
    ("Colaboradores não aparecem", "Verifique se a SOC do colaborador coincide com a SOC do seu perfil"),
    ("Prova não aparece no treinamento", "Verifique se as questões foram cadastradas em Configurações → Avaliações"),
    ("Líder não vê seu time", "Verifique se a 'Chave de Identificação' está EXATAMENTE igual à coluna 'Líder' da planilha"),
    ("Sincronização não roda", "Verifique o horário configurado e se o link do Google Sheets está publicado"),
    ("Vídeo não carrega no treinamento", "Verifique se o link do Google Drive é de COMPARTILHAMENTO e não de edição"),
    ("QR Code não gera", "Verifique se o material tem um link válido cadastrado"),
    ("Solicitação não aparece para análise", "Verifique se você está logado como ADMIN e na aba 'Solicitações'"),
]

for i, (problem, solution) in enumerate(problems):
    y = Inches(1.5 + i * 0.63)
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), y, Inches(12), Inches(0.55))
    s.fill.solid(); s.fill.fore_color.rgb = DARK2; s.line.fill.background()
    add_text(sl, Inches(0.8), y + Inches(0.05), Inches(4.5), Inches(0.22), f"❓ {problem}", 11, ACCENT, True)
    add_text(sl, Inches(0.8), y + Inches(0.28), Inches(11.5), Inches(0.22), f"✅ {solution}", 10, LIGHT)

add_slide_number(sl, slide_num, TOTAL_SLIDES)

# ╔═══════════════════════════════════════════════════════════════╗
# ║  SLIDE 37 — FAQ                                              ║
# ╚═══════════════════════════════════════════════════════════════╝
slide_num += 1
print(f"  [{slide_num:02d}] FAQ")
sl = new_slide()
add_section_header(sl, "❓", "PERGUNTAS FREQUENTES", "Respostas rápidas para dúvidas comuns")

faqs = [
    ("Quantos usuários posso criar?", "Não há limite. Crie quantos forem necessários."),
    ("Posso ter mais de um ADMIN por SOC?", "Sim! Cada SOC pode ter múltiplos administradores."),
    ("O que acontece se eu excluir uma SOC?", "Dados vinculados podem ser afetados. Cuidado ao excluir."),
    ("Posso alterar o e-mail de um usuário?", "Não diretamente. Exclua e recrie o usuário com o novo e-mail."),
    ("A plataforma funciona offline?", "Não. É necessário conexão com internet para acessar."),
    ("Os dados estão seguros?", "Sim. Usamos Supabase/PostgreSQL com banco em nuvem e backup automático."),
    ("Posso acessar pelo celular?", "Sim! A plataforma é responsiva — funciona em qualquer dispositivo."),
    ("Quem faz a sincronização técnica da ABS?", "O desenvolvedor (DEV). Os PTS apenas configuram o horário na plataforma."),
]

for i, (question, answer) in enumerate(faqs):
    y = Inches(1.5 + i * 0.7)
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), y, Inches(12), Inches(0.6))
    s.fill.solid(); s.fill.fore_color.rgb = DARK2; s.line.fill.background()
    add_text(sl, Inches(0.8), y + Inches(0.05), Inches(5), Inches(0.25), f"❓ {question}", 12, ACCENT, True)
    add_text(sl, Inches(6.0), y + Inches(0.05), Inches(6.5), Inches(0.4), f"→ {answer}", 11, LIGHT)

add_slide_number(sl, slide_num, TOTAL_SLIDES)

# ╔═══════════════════════════════════════════════════════════════╗
# ║  SLIDE 38 — GLOSSÁRIO                                       ║
# ╚═══════════════════════════════════════════════════════════════╝
slide_num += 1
print(f"  [{slide_num:02d}] Glossário")
sl = new_slide()
add_section_header(sl, "📖", "GLOSSÁRIO", "Termos e siglas utilizadas na plataforma")

terms = [
    ("SOC", "Sorting & Operational Center — Unidade operacional/hub"),
    ("PTS", "People & Training Specialist — Especialista de Treinamento"),
    ("ABS", "Absence Management System — Sistema de gestão de pessoal"),
    ("BPO", "Business Process Outsourcing — Operação terceirizada"),
    ("PCP", "Planejamento e Controle da Produção"),
    ("EAD", "Ensino a Distância — Treinamento online via plataforma"),
    ("OPSID", "Identificação operacional do colaborador no sistema"),
    ("NR", "Norma Regulamentadora (NR-12, NR-35, etc.)"),
    ("HSE", "Health, Safety & Environment — Saúde, Segurança e Meio Ambiente"),
    ("D+2", "Regra de agendamento: mínimo 2 dias de antecedência"),
    ("QR Code", "Código escaneável para registro rápido de assinatura"),
    ("Badge", "Indicador visual de módulo concluído no onboarding"),
    ("Sync", "Sincronização automática da base de colaboradores"),
    ("Slot", "Horário reservado no calendário para um treinamento"),
]

for i, (term, definition) in enumerate(terms):
    col = i // 7
    row = i % 7
    x = Inches(0.8 + col * 6.3)
    y = Inches(1.5 + row * 0.78)
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.8), Inches(0.68))
    s.fill.solid(); s.fill.fore_color.rgb = DARK2; s.line.fill.background()
    add_text(sl, x + Inches(0.2), y + Inches(0.05), Inches(1.2), Inches(0.3), term, 12, ACCENT, True)
    add_text(sl, x + Inches(1.5), y + Inches(0.05), Inches(4), Inches(0.5), definition, 10, LIGHT)

add_slide_number(sl, slide_num, TOTAL_SLIDES)

# ╔═══════════════════════════════════════════════════════════════╗
# ║  SLIDE 39 — FLUXOGRAMA: PROCESSO COMPLETO                   ║
# ╚═══════════════════════════════════════════════════════════════╝
slide_num += 1
print(f"  [{slide_num:02d}] Fluxograma do Processo")
sl = new_slide()
add_section_header(sl, "🔀", "FLUXO COMPLETO DO ADMINISTRADOR", "Visão macro de todas as responsabilidades do perfil ADMIN")

# Flow boxes
flow = [
    ("1", "CADASTRAR\nSOC", ACCENT, Inches(0.5), Inches(2.0)),
    ("2", "CRIAR\nUSUÁRIOS", ACCENT, Inches(2.5), Inches(2.0)),
    ("3", "CADASTRAR\nINSTRUTORES", ACCENT, Inches(4.5), Inches(2.0)),
    ("4", "SUBIR\nMATERIAIS", ACCENT, Inches(6.5), Inches(2.0)),
    ("5", "CRIAR\nTREINAMENTOS", ACCENT, Inches(8.5), Inches(2.0)),
    ("6", "CONFIGURAR\nPROVAS", ACCENT, Inches(10.5), Inches(2.0)),
    ("7", "IMPORTAR\nCOLABORADORS", GREEN, Inches(0.5), Inches(4.0)),
    ("8", "MONTAR\nAGENDA", GREEN, Inches(2.5), Inches(4.0)),
    ("9", "ANALISAR\nSOLICITAÇÕES", GREEN, Inches(4.5), Inches(4.0)),
    ("10", "INSCREVER\nCOLABORADORS", GREEN, Inches(6.5), Inches(4.0)),
    ("11", "MONITORAR\nRELATÓRIOS", BLUE, Inches(8.5), Inches(4.0)),
    ("12", "AUDITAR\nASSINATURAS", BLUE, Inches(10.5), Inches(4.0)),
]

for num, label, color, x, y in flow:
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(1.7), Inches(1.3))
    s.fill.solid(); s.fill.fore_color.rgb = DARK2; s.line.color.rgb = color; s.line.width = Pt(2)
    # Number badge
    badge = sl.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.6), y - Inches(0.2), Inches(0.45), Inches(0.45))
    badge.fill.solid(); badge.fill.fore_color.rgb = color; badge.line.fill.background()
    add_text(sl, x + Inches(0.68), y - Inches(0.17), Inches(0.3), Inches(0.4), num, 14, WHITE, True, PP_ALIGN.CENTER)
    # Label
    add_text(sl, x + Inches(0.1), y + Inches(0.35), Inches(1.5), Inches(0.8), label, 11, WHITE, True, PP_ALIGN.CENTER)

# Arrows between items (row 1)
for i in range(5):
    x = Inches(2.2 + i * 2)
    s = sl.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, Inches(2.5), Inches(0.3), Inches(0.2))
    s.fill.solid(); s.fill.fore_color.rgb = ACCENT; s.line.fill.background()

# Arrow from row 1 to row 2
s = sl.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(11.2), Inches(3.3), Inches(0.2), Inches(0.5))
s.fill.solid(); s.fill.fore_color.rgb = ACCENT; s.line.fill.background()

# Legend
add_text(sl, Inches(0.8), Inches(5.8), Inches(3), Inches(0.3), "🔴 Configuração Inicial", 12, ACCENT, True)
add_text(sl, Inches(3.8), Inches(5.8), Inches(3), Inches(0.3), "🟢 Operação Diária", 12, GREEN, True)
add_text(sl, Inches(6.8), Inches(5.8), Inches(3), Inches(0.3), "🔵 Monitoramento", 12, BLUE, True)

add_text(sl, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.8),
         "Este fluxo representa a ordem ideal de configuração. A linha superior (vermelha) deve ser feita "
         "uma vez na implantação. A linha inferior (verde/azul) são atividades do dia a dia.", 11, GRAY)

add_slide_number(sl, slide_num, TOTAL_SLIDES)

# ╔═══════════════════════════════════════════════════════════════╗
# ║  SLIDE 40 — SEGURANÇA E ISOLAMENTO                          ║
# ╚═══════════════════════════════════════════════════════════════╝
slide_num += 1
print(f"  [{slide_num:02d}] Segurança e Isolamento")
sl = new_slide()
add_section_header(sl, "🔐", "SEGURANÇA E ISOLAMENTO POR SOC", "Como a plataforma protege os dados de cada unidade")

features = [
    ("🔐", "Isolamento por SOC", "Cada unidade vê APENAS seus próprios dados. Impossível acessar dados de outra SOC."),
    ("🔑", "Senha provisória", "Primeiro login exige redefinição de senha obrigatória."),
    ("👤", "Perfis granulares", "5 níveis de acesso com permissões específicas por módulo."),
    ("☁️", "Banco na nuvem", "Supabase/PostgreSQL com backup automático e criptografia."),
    ("📱", "Multi-dispositivo", "Funciona em desktop, tablet e smartphone — design responsivo."),
    ("📋", "Auditoria completa", "Toda ação é registrada: inscrições, exclusões, aprovações."),
    ("🛡️", "Row Level Security", "Políticas de segurança no banco de dados que impedem acesso indevido."),
    ("🌐", "HTTPS obrigatório", "Toda comunicação é criptografada via protocolo HTTPS."),
]

for i, (icon, title, desc) in enumerate(features):
    y = Inches(1.5 + i * 0.7)
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(11.5), Inches(0.6))
    s.fill.solid(); s.fill.fore_color.rgb = DARK2; s.line.fill.background()
    add_text(sl, Inches(1.0), y + Inches(0.08), Inches(0.4), Inches(0.4), icon, 18, WHITE)
    add_text(sl, Inches(1.6), y + Inches(0.08), Inches(3), Inches(0.35), title, 13, WHITE, True)
    add_text(sl, Inches(4.8), y + Inches(0.08), Inches(7.3), Inches(0.35), desc, 11, GRAY)

add_slide_number(sl, slide_num, TOTAL_SLIDES)

# ╔═══════════════════════════════════════════════════════════════╗
# ║  SLIDE 41 — CONTATOS E SUPORTE                              ║
# ╚═══════════════════════════════════════════════════════════════╝
slide_num += 1
print(f"  [{slide_num:02d}] Contatos e Suporte")
sl = new_slide()
add_section_header(sl, "📞", "CONTATOS E SUPORTE", "Quem procurar em caso de dúvidas ou problemas")

contacts = [
    ("🔧", "SUPORTE TÉCNICO (DEV)", "Problemas no sistema, bugs, sincronização ABS", "Desenvolvedor responsável"),
    ("📋", "PTS PRINCIPAL", "Dúvidas sobre uso da plataforma, configurações", "PTS Admin da SOC principal"),
    ("🏢", "PTS LOCAL", "Configuração da SOC, materiais específicos", "PTS Admin de cada unidade"),
]

for i, (icon, title, desc, who) in enumerate(contacts):
    y = Inches(1.5 + i * 1.5)
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(11.5), Inches(1.3))
    s.fill.solid(); s.fill.fore_color.rgb = DARK2; s.line.color.rgb = BORDER
    add_text(sl, Inches(1.0), y + Inches(0.15), Inches(0.5), Inches(0.5), icon, 28, ACCENT)
    add_text(sl, Inches(1.8), y + Inches(0.15), Inches(5), Inches(0.35), title, 16, WHITE, True)
    add_text(sl, Inches(1.8), y + Inches(0.55), Inches(9), Inches(0.3), desc, 12, LIGHT)
    add_text(sl, Inches(1.8), y + Inches(0.85), Inches(9), Inches(0.3), f"Responsável: {who}", 11, GRAY)

add_info_card(sl, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.8),
              "CANAL DE COMUNICAÇÃO", [
                  "Para problemas técnicos: abra um chamado descrevendo o erro, print da tela e passos para reproduzir"
              ])

add_slide_number(sl, slide_num, TOTAL_SLIDES)

# ╔═══════════════════════════════════════════════════════════════╗
# ║  SLIDE 42 — ENCERRAMENTO                                    ║
# ╚═══════════════════════════════════════════════════════════════╝
slide_num += 1
print(f"  [{slide_num:02d}] Encerramento")
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)
add_accent_bar(sl, Inches(7.44), Inches(0.06))

s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.5), W, Inches(4.5))
s.fill.solid(); s.fill.fore_color.rgb = DARK2; s.line.fill.background()

add_text(sl, Inches(1), Inches(2), Inches(11), Inches(1), "MATRIX ASCEND", 60, ACCENT, True, PP_ALIGN.CENTER)
add_text(sl, Inches(1), Inches(3.0), Inches(11), Inches(0.6),
         "Guia Completo do Administrador", 24, WHITE, True, PP_ALIGN.CENTER)
add_text(sl, Inches(1), Inches(3.6), Inches(11), Inches(0.6),
         "Agora você está preparado para configurar e gerenciar", 18, LIGHT, False, PP_ALIGN.CENTER)
add_text(sl, Inches(1), Inches(4.1), Inches(11), Inches(0.6),
         "a plataforma Matrix Ascend na sua SOC.", 18, LIGHT, False, PP_ALIGN.CENTER)

s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(4.8), Inches(2.3), Inches(0.04))
s.fill.solid(); s.fill.fore_color.rgb = ACCENT; s.line.fill.background()

add_text(sl, Inches(1), Inches(5.2), Inches(11), Inches(0.5),
         "Shopee Logistics  •  SPX BR  •  2026", 14, GRAY, False, PP_ALIGN.CENTER)
add_text(sl, Inches(1), Inches(5.7), Inches(11), Inches(0.4),
         "Bom treinamento! 🎯", 28, WHITE, True, PP_ALIGN.CENTER)

add_slide_number(sl, slide_num, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════
assert slide_num == TOTAL_SLIDES, f"Expected {TOTAL_SLIDES} slides but got {slide_num}"
prs.save(OUT)
print()
print(f"✅ PowerPoint salvo com sucesso!")
print(f"   📄 {OUT}")
print(f"   📊 Total: {slide_num} slides")
print(f"   📐 Formato: 13.333\" × 7.5\" (Widescreen)")
print()
print("🎯 Conteúdo do treinamento:")
print("   • Visão geral da plataforma e navegação")
print("   • Login e primeiro acesso")
print("   • Dashboard e KPIs")
print("   • Configurações: Usuários, Instrutores, Avaliações, Sync")
print("   • Materiais: Pastas, Links, QR Code")
print("   • Colaboradores: CRUD, CSV, Sync, Badges")
print("   • Treinamentos EAD: Vídeo, Prova, Assinatura")
print("   • Agenda: Calendário, Slots, Solicitações, Aprovações")
print("   • SOCs: Cadastro e gestão")
print("   • Assinaturas e Relatórios")
print("   • Permissões, Checklist, FAQ, Glossário")
print("   • Segurança, Troubleshooting, Suporte")
