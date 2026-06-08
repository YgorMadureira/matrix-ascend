from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

BRAIN = r"C:\Users\SPXBR38903\.gemini\antigravity\brain\da637c84-d029-4a63-a983-fe93d2312b58"
OUT = r"C:\Users\SPXBR38903\Documents\Projetos AI\Matrix_Ascend_Apresentacao.pptx"

DARK = RGBColor(0x0D, 0x11, 0x17)
DARK2 = RGBColor(0x16, 0x1B, 0x22)
ACCENT = RGBColor(0xEE, 0x4D, 0x2D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x8B, 0x94, 0x9E)
LIGHT = RGBColor(0xC9, 0xD1, 0xD9)
GREEN = RGBColor(0x3F, 0xB9, 0x50)

SCREENSHOTS = {
    "dashboard": os.path.join(BRAIN, "real_dashboard.png"),
    "collaborators": os.path.join(BRAIN, "real_collaborators.png"),
    "reports": os.path.join(BRAIN, "real_reports.png"),
    "schedule": os.path.join(BRAIN, "real_schedule_requests.png"),
    "signatures": os.path.join(BRAIN, "real_signatures.png"),
    "materials": os.path.join(BRAIN, "real_materials.png"),
    "trainings": os.path.join(BRAIN, "real_trainings.png"),
    "settings": os.path.join(BRAIN, "real_settings.png"),
}

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

def dark_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK

def add_accent_bar(slide, top=0, height=Inches(0.06)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, top, W, height)
    s.fill.solid(); s.fill.fore_color.rgb = ACCENT; s.line.fill.background()

def add_text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.alignment = align
    r = p.runs[0]; r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold; r.font.name = font_name
    return txBox

def add_img(slide, path, left, top, width, height=None):
    if os.path.exists(path):
        if height:
            slide.shapes.add_picture(path, left, top, width, height)
        else:
            slide.shapes.add_picture(path, left, top, width=width)

# ═══════════════════════════════════════
# SLIDE 1 — CAPA
# ═══════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)
add_accent_bar(sl)

# Gradient overlay shape
s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1), W, Inches(5.5))
s.fill.solid(); s.fill.fore_color.rgb = DARK2; s.line.fill.background()

add_text(sl, Inches(1), Inches(1.8), Inches(11), Inches(1), "MATRIX ASCEND", 54, ACCENT, True)
add_text(sl, Inches(1), Inches(2.8), Inches(10), Inches(0.8), "Plataforma Inteligente de Gestão de Treinamentos", 28, WHITE, False)
add_text(sl, Inches(1), Inches(3.6), Inches(10), Inches(0.5), "Shopee Logistics  •  SPX BR", 16, GRAY)

# Linha separadora
s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(4.3), Inches(2), Inches(0.04))
s.fill.solid(); s.fill.fore_color.rgb = ACCENT; s.line.fill.background()

add_text(sl, Inches(1), Inches(4.6), Inches(8), Inches(0.5), "Documento Executivo — Apresentação à Diretoria", 14, GRAY)
add_text(sl, Inches(1), Inches(6.5), Inches(4), Inches(0.4), "v2.4  •  Sistema Online e Operacional", 12, GREEN, True)
add_text(sl, Inches(8), Inches(6.5), Inches(4), Inches(0.4), "© 2026 Matrix Ascend", 12, GRAY, False, PP_ALIGN.RIGHT)

# ═══════════════════════════════════════
# SLIDE 2 — O PROBLEMA
# ═══════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl); add_accent_bar(sl)

add_text(sl, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8), "O PROBLEMA", 36, ACCENT, True)
add_text(sl, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5), "Desafios operacionais que motivaram a criação da plataforma", 14, GRAY)

problems = [
    ("📋", "Controle em planilhas manuais", "Dados desatualizados, duplicados e sem rastreabilidade"),
    ("🔍", "Sem visibilidade de pendências", "Impossível saber quem está certificado em tempo real"),
    ("📝", "Assinaturas em papel", "Risco de fraude, extravio e impossibilidade de auditoria"),
    ("👤", "Onboarding sem rastreio", "Sem visão de quais módulos foram concluídos"),
    ("📧", "Agendamento por e-mail/chat", "Sem controle de inscrições e responsáveis"),
    ("🔒", "Líderes sem autonomia", "Dependência total da área de treinamento"),
    ("🏭", "Múltiplas SOCs sem padrão", "Cada unidade operando de forma diferente"),
]

for i, (icon, title, desc) in enumerate(problems):
    y = Inches(1.8 + i * 0.72)
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(11.5), Inches(0.62))
    s.fill.solid(); s.fill.fore_color.rgb = DARK2; s.line.fill.background()
    add_text(sl, Inches(1.0), y + Inches(0.05), Inches(0.5), Inches(0.5), icon, 18, WHITE)
    add_text(sl, Inches(1.6), y + Inches(0.05), Inches(3.5), Inches(0.5), title, 14, WHITE, True)
    add_text(sl, Inches(5.2), y + Inches(0.05), Inches(7), Inches(0.5), desc, 13, GRAY)

# ═══════════════════════════════════════
# SLIDE 3 — A SOLUÇÃO
# ═══════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl); add_accent_bar(sl)

add_text(sl, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8), "A SOLUÇÃO", 36, ACCENT, True)
add_text(sl, Inches(0.8), Inches(1.2), Inches(10), Inches(0.8), "10 módulos integrados que cobrem todo o ciclo de treinamento", 16, LIGHT)

modules = [
    ("📊", "Dashboard", "KPIs em tempo real"),
    ("👥", "Colaboradores", "Base unificada + Onboarding"),
    ("📁", "Materiais", "Repositório + QR Code"),
    ("🎓", "Capacitação EAD", "Vídeo + Prova + Assinatura"),
    ("✍️", "Assinatura QR", "Coleta mobile em campo"),
    ("📋", "Assinaturas", "Auditoria digital completa"),
    ("📈", "Relatórios", "Matrizes + Gráficos"),
    ("📅", "Agenda & Aprovação", "Fluxo PTS + Google Cal"),
    ("🏢", "SOCs", "Gestão multi-unidade"),
    ("⚙️", "Configurações", "Usuários + Avaliações"),
]

for i, (icon, name, desc) in enumerate(modules):
    col = i % 5
    row = i // 5
    x = Inches(0.8 + col * 2.4)
    y = Inches(2.2 + row * 2.2)
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.1), Inches(1.8))
    s.fill.solid(); s.fill.fore_color.rgb = DARK2; s.line.color.rgb = RGBColor(0x30, 0x36, 0x3D)
    add_text(sl, x + Inches(0.2), y + Inches(0.2), Inches(1.7), Inches(0.5), icon, 28, WHITE)
    add_text(sl, x + Inches(0.2), y + Inches(0.8), Inches(1.7), Inches(0.4), name, 13, WHITE, True)
    add_text(sl, x + Inches(0.2), y + Inches(1.2), Inches(1.7), Inches(0.4), desc, 10, GRAY)

# ═══════════════════════════════════════
# SLIDES 4-11 — SCREENSHOTS DOS MÓDULOS
# ═══════════════════════════════════════
module_slides = [
    ("dashboard", "DASHBOARD", "Painel gerencial com KPIs, % de certificação por setor e visão personalizada por perfil"),
    ("collaborators", "COLABORADORES", "Base unificada com Onboarding, badges por módulo, filtros multi-módulo e importação CSV"),
    ("reports", "RELATÓRIOS & MATRIZ", "Gráficos por SOC, Matriz de Onboarding, Certificação Operacional e Ranking de Instrutores"),
    ("schedule", "AGENDA E APROVAÇÃO (PTS)", "Fluxo de solicitação D+2, análise PTS (Aprovar/Recusar), integração Google Calendar e auditoria"),
    ("signatures", "ASSINATURAS REGISTRADAS", "Auditoria digital com imagem, download individual e filtros por SOC"),
    ("materials", "MATERIAIS", "Repositório de pastas com links Google Drive e geração automática de QR Code"),
    ("trainings", "CAPACITAÇÃO EAD", "Vídeo com timer obrigatório, prova certificadora com 90% mínimo e assinatura digital"),
    ("settings", "CONFIGURAÇÕES", "Gestão de usuários, instrutores, banco de questões e sincronização automática"),
]

for key, title, desc in module_slides:
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(sl); add_accent_bar(sl)
    add_text(sl, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7), title, 30, ACCENT, True)
    add_text(sl, Inches(0.8), Inches(0.9), Inches(11), Inches(0.5), desc, 13, GRAY)

    # Screenshot with border
    img_path = SCREENSHOTS.get(key, "")
    if os.path.exists(img_path):
        # Dark border around screenshot
        s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.55), Inches(11.9), Inches(5.7))
        s.fill.solid(); s.fill.fore_color.rgb = RGBColor(0x21, 0x26, 0x2D)
        s.line.color.rgb = RGBColor(0x30, 0x36, 0x3D)
        add_img(sl, img_path, Inches(0.85), Inches(1.7), Inches(11.6), Inches(5.4))

# ═══════════════════════════════════════
# SLIDE 12 — IMPACTOS
# ═══════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl); add_accent_bar(sl)

add_text(sl, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8), "IMPACTOS POSITIVOS", 36, ACCENT, True)

impacts = [
    ("Verificar pendências", "30-60 min", "2 cliques", "~95%"),
    ("Coletar assinaturas", "Papel → digitalizar", "QR Code automático", "100%"),
    ("Gerar relatórios", "Horas em planilhas", "Tempo real", "99%"),
    ("Agendar treinamento", "E-mails + WhatsApp", "Fluxo PTS + Automático", "~90%"),
    ("Onboarding", "Planilha separada", "Badges por módulo", "100%"),
    ("Atualizar base", "Digitação manual", "Sync automático", "100%"),
]

# Header
add_text(sl, Inches(0.8), Inches(1.4), Inches(3.5), Inches(0.4), "PROCESSO", 11, ACCENT, True)
add_text(sl, Inches(4.5), Inches(1.4), Inches(3), Inches(0.4), "ANTES", 11, ACCENT, True)
add_text(sl, Inches(7.5), Inches(1.4), Inches(3), Inches(0.4), "DEPOIS", 11, ACCENT, True)
add_text(sl, Inches(10.8), Inches(1.4), Inches(1.5), Inches(0.4), "GANHO", 11, ACCENT, True)

for i, (proc, before, after, gain) in enumerate(impacts):
    y = Inches(1.9 + i * 0.75)
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), y, Inches(11.9), Inches(0.65))
    s.fill.solid(); s.fill.fore_color.rgb = DARK2; s.line.fill.background()
    add_text(sl, Inches(0.8), y + Inches(0.1), Inches(3.5), Inches(0.4), proc, 13, WHITE, True)
    add_text(sl, Inches(4.5), y + Inches(0.1), Inches(3), Inches(0.4), before, 12, RGBColor(0xF8, 0x5E, 0x5E))
    add_text(sl, Inches(7.5), y + Inches(0.1), Inches(3), Inches(0.4), after, 12, GREEN)
    add_text(sl, Inches(10.8), y + Inches(0.1), Inches(1.5), Inches(0.4), gain, 14, GREEN, True, PP_ALIGN.CENTER)

# ═══════════════════════════════════════
# SLIDE 13 — SEGURANÇA
# ═══════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl); add_accent_bar(sl)

add_text(sl, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8), "SEGURANÇA E GOVERNANÇA", 36, ACCENT, True)

roles = [
    ("Admin", "Todos os módulos", "Acesso irrestrito"),
    ("Usuário", "Dashboard, Materiais, Colaboradores, Relatórios, SOCs, Treinamentos, Assinaturas", "Sem Configurações"),
    ("Líder", "Dashboard, Colaboradores, Relatórios, Treinamentos, Agenda", "Vê apenas seu time"),
    ("BPO", "Colaboradores (Onboarding)", "Não deleta registros"),
    ("PCP", "Agenda", "Gerencia inscrições"),
]

add_text(sl, Inches(0.8), Inches(1.3), Inches(3), Inches(0.4), "PERFIL", 11, ACCENT, True)
add_text(sl, Inches(3.2), Inches(1.3), Inches(6), Inches(0.4), "MÓDULOS PERMITIDOS", 11, ACCENT, True)
add_text(sl, Inches(9.5), Inches(1.3), Inches(3), Inches(0.4), "RESTRIÇÃO", 11, ACCENT, True)

for i, (role, modules_text, restriction) in enumerate(roles):
    y = Inches(1.8 + i * 0.7)
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), y, Inches(11.9), Inches(0.6))
    s.fill.solid(); s.fill.fore_color.rgb = DARK2; s.line.fill.background()
    add_text(sl, Inches(0.8), y + Inches(0.08), Inches(2.2), Inches(0.4), role, 14, WHITE, True)
    add_text(sl, Inches(3.2), y + Inches(0.08), Inches(6), Inches(0.4), modules_text, 11, GRAY)
    add_text(sl, Inches(9.5), y + Inches(0.08), Inches(3), Inches(0.4), restriction, 11, LIGHT)

# Extra security points
sec_points = [
    "🔐  Isolamento por SOC — dados separados por unidade, impossível ver outra SOC",
    "🔑  Senha provisória com redefinição obrigatória no primeiro login",
    "☁️  Banco de dados em nuvem com backup automático (Supabase/PostgreSQL)",
    "📱  Funciona em qualquer dispositivo — desktop, tablet e smartphone",
]
for i, pt in enumerate(sec_points):
    add_text(sl, Inches(0.8), Inches(5.4 + i * 0.45), Inches(11), Inches(0.4), pt, 12, LIGHT)

# ═══════════════════════════════════════
# SLIDE 14 — ARGUMENTOS FINAIS
# ═══════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl); add_accent_bar(sl)

add_text(sl, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8), "POR QUE APROVAR?", 36, ACCENT, True)

args = [
    ("⏱️", "Redução drástica de tempo", "Processos que consumiam horas agora levam segundos"),
    ("📄", "Eliminação total do papel", "100% digital — listas de presença e registros"),
    ("🔍", "Rastreabilidade completa", "Auditoria digital irrefutável para órgãos reguladores"),
    ("📈", "Escalabilidade imediata", "Adicionar novas SOCs é questão de minutos"),
    ("💰", "Custo mínimo de infraestrutura", "Sem servidores dedicados, sem licenças caras"),
    ("🔒", "Segurança de dados", "Isolamento por unidade com perfis granulares"),
    ("✅", "Já está funcionando", "Não é conceito — é plataforma operacional em produção"),
]

for i, (icon, title, desc) in enumerate(args):
    y = Inches(1.4 + i * 0.78)
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(11.5), Inches(0.68))
    s.fill.solid(); s.fill.fore_color.rgb = DARK2; s.line.fill.background()
    add_text(sl, Inches(1.0), y + Inches(0.08), Inches(0.5), Inches(0.5), icon, 20, WHITE)
    add_text(sl, Inches(1.7), y + Inches(0.08), Inches(4), Inches(0.5), title, 15, WHITE, True)
    add_text(sl, Inches(5.8), y + Inches(0.08), Inches(6.3), Inches(0.5), desc, 13, GRAY)

# ═══════════════════════════════════════
# SLIDE 15 — ENCERRAMENTO
# ═══════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)
add_accent_bar(sl, Inches(7.44), Inches(0.06))

s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.5), W, Inches(4.5))
s.fill.solid(); s.fill.fore_color.rgb = DARK2; s.line.fill.background()

add_text(sl, Inches(1), Inches(2), Inches(11), Inches(1), "MATRIX ASCEND", 60, ACCENT, True, PP_ALIGN.CENTER)
add_text(sl, Inches(1), Inches(3.2), Inches(11), Inches(0.6), "A pessoa certa, treinada na certificação certa,", 22, WHITE, False, PP_ALIGN.CENTER)
add_text(sl, Inches(1), Inches(3.8), Inches(11), Inches(0.6), "com a prova digital irrefutável de que isso aconteceu.", 22, WHITE, False, PP_ALIGN.CENTER)

s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(4.8), Inches(2.3), Inches(0.04))
s.fill.solid(); s.fill.fore_color.rgb = ACCENT; s.line.fill.background()

add_text(sl, Inches(1), Inches(5.2), Inches(11), Inches(0.5), "Shopee Logistics  •  SPX BR  •  2026", 14, GRAY, False, PP_ALIGN.CENTER)
add_text(sl, Inches(1), Inches(5.7), Inches(11), Inches(0.4), "Obrigado.", 28, WHITE, True, PP_ALIGN.CENTER)

# SAVE
prs.save(OUT)
print(f"PowerPoint salvo em: {OUT}")
